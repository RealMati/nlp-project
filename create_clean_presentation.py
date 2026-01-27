#!/usr/bin/env python3
"""
TradeBridge Presentation - Clean with Corner Decorations
Geometric shapes and blobs at corners only
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
import random

# Colors
DARK_BLUE = RGBColor(30, 58, 138)
PRIMARY_BLUE = RGBColor(59, 130, 246)
LIGHT_BLUE = RGBColor(191, 219, 254)
TEAL = RGBColor(20, 184, 166)
GREEN = RGBColor(34, 197, 94)
AMBER = RGBColor(245, 158, 11)
PURPLE = RGBColor(139, 92, 246)
DARK = RGBColor(30, 41, 59)
GRAY = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)

# Decoration colors (softer versions)
DECOR_COLORS = [
    RGBColor(191, 219, 254),  # Light blue
    RGBColor(167, 243, 208),  # Light green
    RGBColor(253, 230, 138),  # Light amber
    RGBColor(196, 181, 253),  # Light purple
    RGBColor(153, 246, 228),  # Light teal
]

def add_corner_decorations(slide, style=0):
    """Add geometric shapes at corners - far from content area"""
    colors = DECOR_COLORS.copy()
    random.shuffle(colors)

    if style == 0:
        # Top-right: circle + small diamond
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.8), Inches(-1), Inches(2), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[0]
        shape.line.fill.background()

        shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(9.5), Inches(0.8), Inches(0.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[1]
        shape.line.fill.background()

        # Bottom-left: circle
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(6.2), Inches(2), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[2]
        shape.line.fill.background()

    elif style == 1:
        # Top-right: triangle
        shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(8.5), Inches(-1.5), Inches(2.5), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[0]
        shape.line.fill.background()

        # Bottom-left: circles cluster
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.3), Inches(5.8), Inches(2.2), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[1]
        shape.line.fill.background()

        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(6.8), Inches(0.8), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[2]
        shape.line.fill.background()

    elif style == 2:
        # Bottom-right: large circle + diamond
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.2), Inches(5.5), Inches(2.5), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[0]
        shape.line.fill.background()

        shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(7.8), Inches(6.5), Inches(0.7), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[1]
        shape.line.fill.background()

        # Top-right: small circle
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(-0.8), Inches(1.2), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[2]
        shape.line.fill.background()

    elif style == 3:
        # Bottom corners only - safer
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(6), Inches(2.5), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[0]
        shape.line.fill.background()

        shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(8.8), Inches(5.8), Inches(2), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[1]
        shape.line.fill.background()

        # Small diamond accent
        shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(8.5), Inches(6.8), Inches(0.5), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[2]
        shape.line.fill.background()

    elif style == 4:
        # Right side decorations only
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(1.8), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[0]
        shape.line.fill.background()

        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(6), Inches(2.2), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[1]
        shape.line.fill.background()

        shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(9.3), Inches(5.5), Inches(0.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[2]
        shape.line.fill.background()

    elif style == 5:
        # Mixed geometric - bottom corners
        shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(-1.2), Inches(5.5), Inches(2.2), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[0]
        shape.line.fill.background()

        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(5.8), Inches(2.2), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[1]
        shape.line.fill.background()

        shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(0.2), Inches(6.8), Inches(0.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[2]
        shape.line.fill.background()

        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.3), Inches(-0.8), Inches(1), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[3]
        shape.line.fill.background()

def add_title(slide, text):
    """Simple title - text only"""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE

    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
    tf = title.text_frame
    tf.text = "TradeBridge"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.6))
    tf = sub.text_frame
    tf.text = "Smart Supply-Demand Management System"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_BLUE

    inst = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
    tf = inst.text_frame
    tf.text = "Adama Science and Technology University\nCollege of Electrical Engineering and Computing"
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.color.rgb = WHITE

    team = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.4))
    tf = team.text_frame
    tf.text = "Hidaya Nurmeka | Ebisa Gutema | Hana Kebede | Hana Jote | Ilham Mohammedhassen"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_BLUE

    adv = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.3))
    tf = adv.text_frame
    tf.text = "Advisor: Dr. Ejigu Tefere"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(245, 158, 11)

def slide_02_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=0)
    add_title(slide, "Agenda")

    items = [
        "01  Problem Overview",
        "02  System Objectives",
        "03  Key Features",
        "04  System Architecture",
        "05  Technology Stack",
        "06  User Roles & Database",
        "07  Machine Learning",
        "08  Benefits & Impact",
        "09  Feasibility & Timeline",
        "10  Conclusion",
    ]

    y = Inches(1.5)
    for i, item in enumerate(items):
        col = i // 5
        row = i % 5
        x = Inches(0.8) + col * Inches(4.5)
        item_y = y + row * Inches(0.9)

        box = slide.shapes.add_textbox(x, item_y, Inches(4), Inches(0.5))
        tf = box.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = DARK

def slide_03_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=1)
    add_title(slide, "Supply Chain Challenges in Ethiopia")

    problems = [
        "- Manual, time-consuming procurement processes",
        "- Lack of centralized B2B marketplace",
        "- Poor demand planning and forecasting",
        "- Limited supplier visibility and comparison",
        "- Inefficient communication between stakeholders",
        "- Stock shortages and delivery delays"
    ]

    y = Inches(1.6)
    for prob in problems:
        box = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.5), Inches(0.5))
        tf = box.text_frame
        tf.text = prob
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = DARK
        y += Inches(0.8)

def slide_04_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=2)
    add_title(slide, "Introducing TradeBridge")

    desc = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.2))
    tf = desc.text_frame
    tf.word_wrap = True
    tf.text = "A comprehensive B2B digital platform connecting retailers, factories, distributors, and delivery personnel to streamline bulk ordering, improve supply chain transparency, and enable data-driven decision-making."
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = DARK

    stakeholders = [
        "Retailers: Browse, compare, order products",
        "Factories: Manage products and orders",
        "Distributors: Inventory management",
        "Delivery: Track routes and status",
    ]

    y = Inches(3.2)
    for item in stakeholders:
        box = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.5), Inches(0.5))
        tf = box.text_frame
        tf.text = f"- {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = DARK
        y += Inches(0.7)

def slide_05_scope(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=3)
    add_title(slide, "Project Scope")

    in_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4), Inches(0.5))
    tf = in_title.text_frame
    tf.text = "IN SCOPE:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(34, 197, 94)

    in_items = [
        "- Food and beverage products",
        "- Micro to large enterprises",
        "- Ethiopian market (ETB)",
        "- Mobile & Web platforms",
        "- Distribution of finished goods"
    ]

    y = Inches(2.1)
    for item in in_items:
        box = slide.shapes.add_textbox(Inches(0.8), y, Inches(4), Inches(0.5))
        tf = box.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        y += Inches(0.5)

    out_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4), Inches(0.5))
    tf = out_title.text_frame
    tf.text = "OUT OF SCOPE:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(239, 68, 68)

    out_items = [
        "- Raw material procurement",
        "- International trade",
        "- Large national producers"
    ]

    y = Inches(2.1)
    for item in out_items:
        box = slide.shapes.add_textbox(Inches(5.2), y, Inches(4), Inches(0.5))
        tf = box.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        y += Inches(0.5)

def slide_06_objectives(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=4)
    add_title(slide, "Key Objectives")

    gen = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(0.8))
    tf = gen.text_frame
    tf.word_wrap = True
    tf.text = "General: Design and develop a digital platform to streamline B2B procurement and enhance supply chain visibility."
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    objectives = [
        "1. Develop Web & Mobile Application",
        "2. Implement ML-based Supplier Recommendation",
        "3. Introduce Demand Forecasting Capabilities",
        "4. Enable Real-time Communication",
        "5. Provide Centralized Supplier Directory"
    ]

    y = Inches(2.5)
    for obj in objectives:
        box = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.5), Inches(0.5))
        tf = box.text_frame
        tf.text = obj
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = DARK
        y += Inches(0.75)

def slide_07_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=5)
    add_title(slide, "Platform Features")

    ret_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(4), Inches(0.5))
    tf = ret_title.text_frame
    tf.text = "For Retailers:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    ret_features = [
        "- Browse and compare products",
        "- Place bulk orders",
        "- Track order status",
        "- Rate suppliers",
        "- Personalized recommendations"
    ]

    y = Inches(2.0)
    for feat in ret_features:
        box = slide.shapes.add_textbox(Inches(0.8), y, Inches(4), Inches(0.5))
        tf = box.text_frame
        tf.text = feat
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        y += Inches(0.5)

    sup_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4), Inches(0.5))
    tf = sup_title.text_frame
    tf.text = "For Suppliers:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(245, 158, 11)

    sup_features = [
        "- Manage product listings",
        "- Approve/reject orders",
        "- Broadcast announcements",
        "- View demand analytics"
    ]

    y = Inches(2.0)
    for feat in sup_features:
        box = slide.shapes.add_textbox(Inches(5.2), y, Inches(4), Inches(0.5))
        tf = box.text_frame
        tf.text = feat
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        y += Inches(0.5)

def slide_08_smart_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=0)
    add_title(slide, "Smart Capabilities")

    features = [
        ("Supplier Recommendation System", "Ranks suppliers based on price, distance, performance using Random Forest"),
        ("Demand Forecasting", "Predicts future demand using historical data and seasonality patterns"),
        ("Real-time GPS Tracking", "Track deliveries in real-time with location updates"),
        ("Chapa Payment Gateway", "Secure payment processing integrated with local banking"),
        ("Auto Notifications", "Automated alerts for orders, deliveries, and updates"),
        ("In-App Messaging", "Direct communication between retailers and suppliers"),
    ]

    y = Inches(1.5)
    for title, desc in features:
        title_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.5), Inches(0.4))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

        desc_box = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.35), Inches(8.5), Inches(0.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = desc
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY

        y += Inches(0.9)

def slide_09_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=1)
    add_title(slide, "Three-Tier Architecture")

    layers = [
        ("PRESENTATION LAYER", "Mobile App | Web Application | User Interface"),
        ("APPLICATION LAYER", "Node.js + Express | Auth | Business Logic | ML Integration"),
        ("DATA LAYER", "MySQL | User Data | Products & Orders | Analytics"),
    ]

    y = Inches(1.8)
    for name, components in layers:
        name_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.5), Inches(0.5))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

        comp_box = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.5), Inches(8.5), Inches(0.5))
        tf = comp_box.text_frame
        tf.text = components
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY

        if name != "DATA LAYER":
            arrow_box = slide.shapes.add_textbox(Inches(4.5), y + Inches(1.1), Inches(1), Inches(0.4))
            tf = arrow_box.text_frame
            tf.text = "v"
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(28)
            p.font.color.rgb = GRAY

        y += Inches(1.6)

def slide_10_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=2)
    add_title(slide, "Technology Stack")

    tech = [
        ("Frontend:", "React.js with TypeScript, Tailwind CSS, Zustand"),
        ("Backend:", "Node.js with Express, JWT Auth, RESTful APIs"),
        ("Database:", "MySQL with Sequelize ORM"),
        ("ML:", "Python, Scikit-learn, Pandas, NumPy"),
        ("Payment:", "Chapa Payment Gateway"),
    ]

    y = Inches(1.6)
    for category, items in tech:
        cat_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(2), Inches(0.5))
        tf = cat_box.text_frame
        tf.text = category
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

        items_box = slide.shapes.add_textbox(Inches(2.8), y, Inches(6.5), Inches(0.5))
        tf = items_box.text_frame
        tf.text = items
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = DARK

        y += Inches(0.9)

def slide_11_user_roles(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=3)
    add_title(slide, "User Roles")

    roles = [
        ("Retailer:", "Browse, order, track, rate suppliers"),
        ("Factory:", "Manage products, orders, view demand forecasts"),
        ("Distributor:", "Handle inventory, manage fulfillment"),
        ("Driver:", "View assignments, update delivery status, navigate routes"),
        ("Admin:", "Manage users, suppliers, system monitoring"),
    ]

    y = Inches(1.6)
    for role, permissions in roles:
        role_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(2.2), Inches(0.5))
        tf = role_box.text_frame
        tf.text = role
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

        perm_box = slide.shapes.add_textbox(Inches(3), y, Inches(6.3), Inches(0.5))
        tf = perm_box.text_frame
        tf.text = permissions
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = DARK

        y += Inches(0.9)

def slide_12_database(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=4)
    add_title(slide, "Database Entities")

    entities = [
        ("Users:", "Retailers, suppliers, drivers, admins"),
        ("Products:", "Name, price, stock, MOQ, images"),
        ("Orders:", "Status, items, tracking, payments"),
        ("Messages:", "In-app chat history"),
        ("Ratings:", "Supplier performance metrics"),
        ("Payments:", "Transactions, methods, status"),
    ]

    y = Inches(1.6)
    for name, desc in entities:
        name_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(2), Inches(0.5))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

        desc_box = slide.shapes.add_textbox(Inches(2.8), y, Inches(6.5), Inches(0.5))
        tf = desc_box.text_frame
        tf.text = desc
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = DARK

        y += Inches(0.85)

def slide_13_ml(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=5)
    add_title(slide, "ML Pipeline")

    steps = [
        ("1. Collect Data", "Price, delivery times, ratings, location data"),
        ("2. Train Model", "Random Forest Classifier for supplier ranking"),
        ("3. Score & Rank", "Rank suppliers per retailer based on preferences"),
        ("4. Adapt & Learn", "Continuous learning from user preferences"),
    ]

    y = Inches(1.6)
    for title, desc in steps:
        title_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.5), Inches(0.5))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

        desc_box = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.45), Inches(8.5), Inches(0.5))
        tf = desc_box.text_frame
        tf.text = desc
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY

        y += Inches(1.1)

    feat_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.6))
    tf = feat_box.text_frame
    tf.word_wrap = True
    tf.text = "Features: Price competitiveness, Delivery rate, Quality ratings, Fulfillment time"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK

def slide_14_benefits(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=0)
    add_title(slide, "Expected Impact")

    benefits = [
        ("For Retailers:", "Faster procurement, better price comparison, improved visibility"),
        ("For Suppliers:", "Expanded reach, direct connections, reduced manual work, demand insights"),
        ("For Industry:", "Digital transformation, reduced inefficiency, increased transparency"),
    ]

    y = Inches(1.8)
    for title, items in benefits:
        title_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.5), Inches(0.5))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

        items_box = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.5), Inches(8.5), Inches(0.6))
        tf = items_box.text_frame
        tf.word_wrap = True
        tf.text = items
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = DARK

        y += Inches(1.5)

def slide_15_feasibility(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=1)
    add_title(slide, "Project Feasibility")

    feasibility = [
        ("Technical:", "Proven technologies, team has skills, no special hardware"),
        ("Operational:", "User-friendly design, supports business practices, mobile-first"),
        ("Economic:", "Low cost (~9,500 ETB), open-source tech, transaction fee revenue"),
    ]

    y = Inches(1.8)
    for title, items in feasibility:
        title_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.5), Inches(0.5))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

        items_box = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.5), Inches(8.5), Inches(0.6))
        tf = items_box.text_frame
        tf.word_wrap = True
        tf.text = items
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = DARK

        y += Inches(1.5)

def slide_16_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=2)
    add_title(slide, "Development Timeline")

    phases = [
        ("Planning:", "Week 1-2"),
        ("Analysis:", "Week 3-8"),
        ("Design:", "Week 8-12"),
        ("Implementation:", "Week 13-15"),
        ("Testing:", "Week 15-16"),
        ("Documentation:", "Week 17-18"),
    ]

    y = Inches(1.6)
    for phase, timeline in phases:
        phase_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(3), Inches(0.5))
        tf = phase_box.text_frame
        tf.text = phase
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

        time_box = slide.shapes.add_textbox(Inches(3.8), y, Inches(5.5), Inches(0.5))
        tf = time_box.text_frame
        tf.text = timeline
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = DARK

        y += Inches(0.75)

    total = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
    tf = total.text_frame
    tf.text = "Total: ~18 weeks (2 semesters)"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

def slide_17_challenges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=3)
    add_title(slide, "Challenges & Solutions")

    ch_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(0.5))
    tf = ch_title.text_frame
    tf.text = "Challenges:"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    challenges = [
        "- Limited initial ML training data",
        "- Internet connectivity dependency",
        "- Digital adoption resistance",
        "- Third-party payment reliance",
    ]

    y = Inches(2.0)
    for ch in challenges:
        box = slide.shapes.add_textbox(Inches(1), y, Inches(8), Inches(0.4))
        tf = box.text_frame
        tf.text = ch
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        y += Inches(0.5)

    mit_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.5), Inches(0.5))
    tf = mit_title.text_frame
    tf.text = "Mitigation Strategies:"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    mitigations = [
        "- Hybrid ML approach (rule-based + data-driven)",
        "- Offline capabilities for critical functions",
        "- Comprehensive user onboarding and support"
    ]

    y = Inches(4.7)
    for mit in mitigations:
        box = slide.shapes.add_textbox(Inches(1), y, Inches(8.2), Inches(0.4))
        tf = box.text_frame
        tf.text = mit
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        y += Inches(0.5)

def slide_18_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_corner_decorations(slide, style=4)
    add_title(slide, "Future Roadmap")

    enhancements = [
        "- Additional product categories",
        "- AI-powered chatbot support",
        "- iOS mobile application",
        "- Multi-language support",
        "- Advanced analytics dashboards",
        "- ERP system integration",
        "- Logistics optimization",
        "- Blockchain transparency",
    ]

    y = Inches(1.6)
    for enhancement in enhancements:
        box = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.5), Inches(0.5))
        tf = box.text_frame
        tf.text = enhancement
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = DARK
        y += Inches(0.65)

def slide_19_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
    tf = title.text_frame
    tf.text = "Summary"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(0.5))
    tf = subtitle.text_frame
    tf.text = "TradeBridge Delivers:"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    deliverables = [
        "- Centralized B2B marketplace for Ethiopian wholesale",
        "- Smart supplier recommendations & demand forecasting",
        "- Real-time tracking & secure payments",
        "- Improved efficiency across the supply chain"
    ]

    y = Inches(2.5)
    for item in deliverables:
        box = slide.shapes.add_textbox(Inches(1), y, Inches(8), Inches(0.5))
        tf = box.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        y += Inches(0.7)

    status_text = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.6))
    tf = status_text.text_frame
    tf.text = "Prototype Completed - Ready for Deployment"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(34, 197, 94)

def slide_20_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = TEAL

    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
    tf = thanks.text_frame
    tf.text = "Thank You"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE

    questions = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.6))
    tf = questions.text_frame
    tf.text = "Questions?"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE

    contact = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.8))
    tf = contact.text_frame
    tf.text = "Advisor: Dr. Ejigu Tefere\nAdama Science and Technology University"
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = WHITE

    team = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.4))
    tf = team.text_frame
    tf.text = "Hidaya Nurmeka | Ebisa Gutema | Hana Kebede | Hana Jote | Ilham Mohammedhassen"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("\nCreating TEXT-ONLY Presentation...\n")

    slides = [
        ("Title", slide_01_title),
        ("Agenda", slide_02_agenda),
        ("Problem", slide_03_problem),
        ("Introduction", slide_04_intro),
        ("Scope", slide_05_scope),
        ("Objectives", slide_06_objectives),
        ("Features", slide_07_features),
        ("Smart Features", slide_08_smart_features),
        ("Architecture", slide_09_architecture),
        ("Tech Stack", slide_10_tech_stack),
        ("User Roles", slide_11_user_roles),
        ("Database", slide_12_database),
        ("ML Pipeline", slide_13_ml),
        ("Benefits", slide_14_benefits),
        ("Feasibility", slide_15_feasibility),
        ("Timeline", slide_16_timeline),
        ("Challenges", slide_17_challenges),
        ("Future", slide_18_future),
        ("Conclusion", slide_19_conclusion),
        ("Thank You", slide_20_thankyou),
    ]

    for name, func in slides:
        print(f"  {name}")
        func(prs)

    output = 'TradeBridge_Clean.pptx'
    prs.save(output)

    print(f"\nDone! File: {os.path.abspath(output)}")
    print(f"Slides: {len(prs.slides)}")

if __name__ == '__main__':
    main()
