#!/usr/bin/env python3
"""
TradeBridge Stunning Presentation - Professional Design with Real Visuals
No emojis, actual shapes/graphics, curved backgrounds, modern corporate style
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Premium Color Palette
DEEP_BLUE = RGBColor(15, 23, 42)
PRIMARY_BLUE = RGBColor(37, 99, 235)
LIGHT_BLUE = RGBColor(96, 165, 250)
PALE_BLUE = RGBColor(219, 234, 254)
TEAL = RGBColor(20, 184, 166)
GREEN = RGBColor(34, 197, 94)
DARK_GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(245, 158, 11)
ORANGE = RGBColor(249, 115, 22)
PURPLE = RGBColor(139, 92, 246)
ROSE = RGBColor(244, 63, 94)
SLATE = RGBColor(71, 85, 105)
DARK_SLATE = RGBColor(30, 41, 59)
LIGHT_SLATE = RGBColor(148, 163, 184)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(248, 250, 252)

def add_curved_accent(slide, x, y, width, height, color, rotation=0):
    """Add curved accent shape"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.WAVE,
        x, y, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = rotation
    return shape

def add_circle_accent(slide, x, y, size, color, transparency=0):
    """Add decorative circle"""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x, y, size, size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    return circle

def add_rounded_card(slide, x, y, w, h, fill_color, border_color=None, border_width=0):
    """Create a card with rounded corners"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, y, w, h
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(border_width)
    else:
        card.line.fill.background()
    return card

def add_icon_shape(slide, x, y, size, icon_type, color):
    """Create icon using shapes"""
    if icon_type == "chart":
        # Bar chart icon
        bar_w = size / 4
        bar_gap = size / 8
        heights = [size * 0.5, size * 0.8, size * 0.6, size * 0.9]
        for i, h in enumerate(heights):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x + i * (bar_w + bar_gap), y + (size - h),
                bar_w, h
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()
    elif icon_type == "gear":
        gear = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, x, y, size, size)
        gear.fill.solid()
        gear.fill.fore_color.rgb = color
        gear.line.fill.background()
    elif icon_type == "lightning":
        bolt = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size, size)
        bolt.fill.solid()
        bolt.fill.fore_color.rgb = color
        bolt.line.fill.background()
    elif icon_type == "arrow_right":
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, size, size * 0.6)
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()
    elif icon_type == "checkmark":
        # Create checkmark using lines/shapes
        check = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size, size)
        check.fill.solid()
        check.fill.fore_color.rgb = color
        check.line.fill.background()
    elif icon_type == "database":
        # Cylinder for database
        db = slide.shapes.add_shape(MSO_SHAPE.CAN, x, y, size, size * 1.2)
        db.fill.solid()
        db.fill.fore_color.rgb = color
        db.line.fill.background()
    elif icon_type == "cloud":
        cloud = slide.shapes.add_shape(MSO_SHAPE.CLOUD, x, y, size, size * 0.7)
        cloud.fill.solid()
        cloud.fill.fore_color.rgb = color
        cloud.line.fill.background()
    elif icon_type == "hexagon":
        hex_shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, x, y, size, size)
        hex_shape.fill.solid()
        hex_shape.fill.fore_color.rgb = color
        hex_shape.line.fill.background()
    elif icon_type == "star":
        star = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, size, size)
        star.fill.solid()
        star.fill.fore_color.rgb = color
        star.line.fill.background()
    elif icon_type == "diamond":
        diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, size, size)
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = color
        diamond.line.fill.background()

def add_decorative_background(slide, style="curves"):
    """Add decorative background elements"""
    if style == "curves":
        # Large curved shape at bottom right
        wave1 = slide.shapes.add_shape(
            MSO_SHAPE.WAVE,
            Inches(5), Inches(4),
            Inches(6), Inches(4)
        )
        wave1.fill.solid()
        wave1.fill.fore_color.rgb = PALE_BLUE
        wave1.line.fill.background()
        wave1.rotation = 15

        # Smaller curve top left
        wave2 = slide.shapes.add_shape(
            MSO_SHAPE.WAVE,
            Inches(-1), Inches(-0.5),
            Inches(4), Inches(3)
        )
        wave2.fill.solid()
        wave2.fill.fore_color.rgb = PALE_BLUE
        wave2.line.fill.background()
        wave2.rotation = -30

    elif style == "circles":
        # Decorative circles
        circles_data = [
            (Inches(8), Inches(-0.5), Inches(2), PALE_BLUE),
            (Inches(9), Inches(5.5), Inches(3), PALE_BLUE),
            (Inches(-0.8), Inches(5), Inches(2.5), PALE_BLUE),
        ]
        for x, y, size, color in circles_data:
            add_circle_accent(slide, x, y, size, color)

    elif style == "geometric":
        # Geometric pattern
        hex1 = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(8.5), Inches(-0.3), Inches(2), Inches(2))
        hex1.fill.solid()
        hex1.fill.fore_color.rgb = PALE_BLUE
        hex1.line.fill.background()
        hex1.rotation = 30

        hex2 = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(-0.5), Inches(5.5), Inches(1.5), Inches(1.5))
        hex2.fill.solid()
        hex2.fill.fore_color.rgb = PALE_BLUE
        hex2.line.fill.background()

def add_title_bar(slide, title, subtitle=None):
    """Add stylish title bar"""
    # Title background strip
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.3)
    )
    strip.fill.gradient()
    strip.fill.gradient_angle = 0
    strip.fill.gradient_stops[0].color.rgb = PRIMARY_BLUE
    strip.fill.gradient_stops[1].color.rgb = TEAL
    strip.line.fill.background()

    # Accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.3),
        Inches(10), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = AMBER
    accent.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(8.4), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(8.4), Inches(0.4))
        tf = sub_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(200, 220, 255)

def slide_01_title(prs):
    """Title Slide - Hero style"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Gradient background
    bg = slide.background
    bg.fill.gradient()
    bg.fill.gradient_angle = 135
    bg.fill.gradient_stops[0].color.rgb = DEEP_BLUE
    bg.fill.gradient_stops[1].color.rgb = PRIMARY_BLUE

    # Decorative curves
    wave1 = slide.shapes.add_shape(MSO_SHAPE.WAVE, Inches(5), Inches(4.5), Inches(7), Inches(4))
    wave1.fill.solid()
    wave1.fill.fore_color.rgb = RGBColor(30, 64, 175)
    wave1.line.fill.background()
    wave1.rotation = 10

    wave2 = slide.shapes.add_shape(MSO_SHAPE.WAVE, Inches(-2), Inches(-1), Inches(5), Inches(3.5))
    wave2.fill.solid()
    wave2.fill.fore_color.rgb = RGBColor(30, 64, 175)
    wave2.line.fill.background()
    wave2.rotation = -20

    # Decorative circles
    add_circle_accent(slide, Inches(8.5), Inches(0.5), Inches(1.2), RGBColor(59, 130, 246))
    add_circle_accent(slide, Inches(0.5), Inches(5.5), Inches(0.8), RGBColor(59, 130, 246))
    add_circle_accent(slide, Inches(9), Inches(5), Inches(0.5), TEAL)

    # Main title
    title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
    tf = title.text_frame
    tf.text = "TradeBridge"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Accent line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.2), Inches(3), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = AMBER
    line.line.fill.background()

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    tf = subtitle.text_frame
    tf.text = "Smart Supply-Demand Management System"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_BLUE

    # Institution card
    inst_card = add_rounded_card(slide, Inches(2), Inches(4.6), Inches(6), Inches(1.1), RGBColor(30, 58, 138))
    tf = inst_card.text_frame
    tf.word_wrap = True
    tf.text = "Adama Science and Technology University\nCollege of Electrical Engineering and Computing"
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(14)
        para.font.color.rgb = WHITE

    # Team
    team = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.4))
    tf = team.text_frame
    tf.text = "Hidaya Nurmeka  |  Ebisa Gutema  |  Hana Kebede  |  Hana Jote  |  Ilham Mohammedhassen"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_SLATE

    # Advisor
    advisor = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.3))
    tf = advisor.text_frame
    tf.text = "Advisor: Dr. Ejigu Tefere"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AMBER

def slide_02_agenda(prs):
    """Agenda slide with modern cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "geometric")
    add_title_bar(slide, "Presentation Agenda")

    agenda_items = [
        ("01", "Problem Overview", PRIMARY_BLUE),
        ("02", "System Objectives", TEAL),
        ("03", "Key Features", GREEN),
        ("04", "System Architecture", PURPLE),
        ("05", "Technology Stack", AMBER),
        ("06", "User Roles & Database", ROSE),
        ("07", "Machine Learning", PRIMARY_BLUE),
        ("08", "Benefits & Impact", TEAL),
        ("09", "Feasibility & Timeline", GREEN),
        ("10", "Conclusion", PURPLE),
    ]

    # Grid layout
    col_width = Inches(4.3)
    row_height = Inches(0.95)
    x_positions = [Inches(0.6), Inches(5.1)]
    y_start = Inches(1.8)

    for i, (num, text, color) in enumerate(agenda_items):
        row = i % 5
        col = i // 5
        x = x_positions[col]
        y = y_start + row * (row_height + Inches(0.15))

        # Card background
        card = add_rounded_card(slide, x, y, col_width, row_height, WHITE, LIGHT_SLATE, 1)

        # Number circle
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.17), Inches(0.6), Inches(0.6))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = color
        num_circle.line.fill.background()

        tf = num_circle.text_frame
        tf.text = num
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Text
        text_box = slide.shapes.add_textbox(x + Inches(0.9), y + Inches(0.25), Inches(3.2), Inches(0.5))
        tf = text_box.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_SLATE

def slide_03_problem(prs):
    """Problem Statement with visual cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "curves")
    add_title_bar(slide, "Supply Chain Challenges in Ethiopia")

    problems = [
        ("Manual, time-consuming procurement processes", MSO_SHAPE.PENTAGON),
        ("Lack of centralized B2B marketplace", MSO_SHAPE.CLOUD),
        ("Poor demand planning and forecasting", MSO_SHAPE.PARALLELOGRAM),
        ("Limited supplier visibility and comparison", MSO_SHAPE.BEVEL),
        ("Inefficient stakeholder communication", MSO_SHAPE.CHEVRON),
        ("Stock shortages and delivery delays", MSO_SHAPE.CUBE),
    ]

    # 2x3 grid
    x_positions = [Inches(0.6), Inches(5.1)]
    y_start = Inches(1.8)
    card_w = Inches(4.3)
    card_h = Inches(1.6)

    for i, (text, icon_shape) in enumerate(problems):
        row = i // 2
        col = i % 2
        x = x_positions[col]
        y = y_start + row * (card_h + Inches(0.2))

        # Card
        card = add_rounded_card(slide, x, y, card_w, card_h, WHITE, ROSE, 3)

        # Icon
        icon = slide.shapes.add_shape(icon_shape, x + Inches(0.3), y + Inches(0.4), Inches(0.7), Inches(0.7))
        icon.fill.solid()
        icon.fill.fore_color.rgb = ROSE
        icon.line.fill.background()

        # Text
        text_box = slide.shapes.add_textbox(x + Inches(1.2), y + Inches(0.35), card_w - Inches(1.5), card_h - Inches(0.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(17)
        p.font.color.rgb = DARK_SLATE

def slide_04_intro(prs):
    """Introduction to TradeBridge"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background
    bg.fill.gradient()
    bg.fill.gradient_angle = 180
    bg.fill.gradient_stops[0].color.rgb = OFF_WHITE
    bg.fill.gradient_stops[1].color.rgb = PALE_BLUE

    add_title_bar(slide, "Introducing TradeBridge")

    # Main description card
    desc_card = add_rounded_card(slide, Inches(0.8), Inches(1.7), Inches(8.4), Inches(1.6), PRIMARY_BLUE)

    # Gradient overlay effect using shape
    desc_overlay = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.7), Inches(8.4), Inches(1.6)
    )
    desc_overlay.fill.gradient()
    desc_overlay.fill.gradient_angle = 90
    desc_overlay.fill.gradient_stops[0].color.rgb = PRIMARY_BLUE
    desc_overlay.fill.gradient_stops[1].color.rgb = TEAL
    desc_overlay.line.fill.background()

    desc_text = slide.shapes.add_textbox(Inches(1.2), Inches(2), Inches(7.6), Inches(1))
    tf = desc_text.text_frame
    tf.word_wrap = True
    tf.text = "A comprehensive B2B digital platform connecting retailers, factories, distributors, and delivery personnel to streamline bulk ordering, improve supply chain transparency, and enable data-driven decision-making."
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE

    # Stakeholder cards
    stakeholders = [
        ("Retailers", "Browse, compare, order", PRIMARY_BLUE, MSO_SHAPE.RECTANGLE),
        ("Factories", "Manage products, forecasts", TEAL, MSO_SHAPE.PENTAGON),
        ("Distributors", "Inventory, fulfillment", GREEN, MSO_SHAPE.CUBE),
        ("Delivery", "Track routes, status", AMBER, MSO_SHAPE.RIGHT_ARROW),
    ]

    x_start = Inches(0.9)
    y = Inches(3.8)
    card_w = Inches(2.1)
    card_h = Inches(2.5)
    gap = Inches(0.25)

    for i, (name, desc, color, icon_shape) in enumerate(stakeholders):
        x = x_start + i * (card_w + gap)

        # Card
        card = add_rounded_card(slide, x, y, card_w, card_h, WHITE, color, 3)

        # Icon shape
        icon = slide.shapes.add_shape(icon_shape, x + Inches(0.65), y + Inches(0.4), Inches(0.8), Inches(0.8))
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.line.fill.background()

        # Name
        name_box = slide.shapes.add_textbox(x, y + Inches(1.35), card_w, Inches(0.4))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color

        # Description
        desc_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.8), card_w - Inches(0.2), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = desc
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE

def slide_05_scope(prs):
    """Project Scope"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "circles")
    add_title_bar(slide, "Project Scope")

    # In Scope Card
    in_card = add_rounded_card(slide, Inches(0.6), Inches(1.7), Inches(4.3), Inches(5), GREEN)
    in_overlay = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.7), Inches(4.3), Inches(5))
    in_overlay.fill.gradient()
    in_overlay.fill.gradient_angle = 180
    in_overlay.fill.gradient_stops[0].color.rgb = GREEN
    in_overlay.fill.gradient_stops[1].color.rgb = DARK_GREEN
    in_overlay.line.fill.background()

    # In Scope header
    in_header = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(3.9), Inches(0.5))
    tf = in_header.text_frame
    tf.text = "IN SCOPE"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Checkmark icon
    check = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(2.35), Inches(2.5), Inches(0.8), Inches(0.8))
    check.fill.solid()
    check.fill.fore_color.rgb = WHITE
    check.line.fill.background()

    in_items = [
        "Food and beverage products",
        "Micro to large-sized enterprises",
        "Ethiopian market (ETB currency)",
        "Mobile & Web platforms",
        "Distribution of finished goods"
    ]

    y = Inches(3.5)
    for item in in_items:
        item_box = slide.shapes.add_textbox(Inches(1), y, Inches(3.5), Inches(0.5))
        tf = item_box.text_frame
        tf.word_wrap = True
        tf.text = f"  {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        y += Inches(0.55)

    # Out of Scope Card
    out_card = add_rounded_card(slide, Inches(5.1), Inches(1.7), Inches(4.3), Inches(5), ROSE)
    out_overlay = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.7), Inches(4.3), Inches(5))
    out_overlay.fill.gradient()
    out_overlay.fill.gradient_angle = 180
    out_overlay.fill.gradient_stops[0].color.rgb = ROSE
    out_overlay.fill.gradient_stops[1].color.rgb = RGBColor(190, 18, 60)
    out_overlay.line.fill.background()

    # Out Scope header
    out_header = slide.shapes.add_textbox(Inches(5.3), Inches(1.9), Inches(3.9), Inches(0.5))
    tf = out_header.text_frame
    tf.text = "OUT OF SCOPE"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # X icon
    x_mark = slide.shapes.add_shape(MSO_SHAPE.CROSS, Inches(6.85), Inches(2.5), Inches(0.8), Inches(0.8))
    x_mark.fill.solid()
    x_mark.fill.fore_color.rgb = WHITE
    x_mark.line.fill.background()

    out_items = [
        "Raw material procurement between factories",
        "International trade operations",
        "Very large national producers (e.g., Wenji Sugar)"
    ]

    y = Inches(3.5)
    for item in out_items:
        item_box = slide.shapes.add_textbox(Inches(5.5), y, Inches(3.5), Inches(0.7))
        tf = item_box.text_frame
        tf.word_wrap = True
        tf.text = f"  {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        y += Inches(0.75)

def slide_06_objectives(prs):
    """System Objectives"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "curves")
    add_title_bar(slide, "Key Objectives")

    # General objective
    gen_card = add_rounded_card(slide, Inches(0.8), Inches(1.6), Inches(8.4), Inches(1.1), PRIMARY_BLUE)
    gen_overlay = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(8.4), Inches(1.1))
    gen_overlay.fill.gradient()
    gen_overlay.fill.gradient_angle = 90
    gen_overlay.fill.gradient_stops[0].color.rgb = PRIMARY_BLUE
    gen_overlay.fill.gradient_stops[1].color.rgb = TEAL
    gen_overlay.line.fill.background()

    gen_text = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(0.7))
    tf = gen_text.text_frame
    tf.word_wrap = True
    tf.text = "Design and develop a digital platform to streamline B2B procurement and enhance supply chain visibility across Ethiopia."
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE

    # Specific objectives
    objectives = [
        ("Develop Web & Mobile Application", MSO_SHAPE.BEVEL, PRIMARY_BLUE),
        ("Implement ML-based Supplier Recommendation", MSO_SHAPE.PENTAGON, TEAL),
        ("Introduce Demand Forecasting Capabilities", MSO_SHAPE.PARALLELOGRAM, GREEN),
        ("Enable Real-time Communication", MSO_SHAPE.CHEVRON, PURPLE),
        ("Provide Centralized Supplier Directory", MSO_SHAPE.CAN, AMBER),
    ]

    y = Inches(3)
    for text, icon_shape, color in objectives:
        # Card
        card = add_rounded_card(slide, Inches(1.2), y, Inches(7.6), Inches(0.8), WHITE, color, 2)

        # Icon
        icon = slide.shapes.add_shape(icon_shape, Inches(1.5), y + Inches(0.12), Inches(0.55), Inches(0.55))
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.line.fill.background()

        # Text
        text_box = slide.shapes.add_textbox(Inches(2.3), y + Inches(0.15), Inches(6.3), Inches(0.5))
        tf = text_box.text_frame
        tf.text = text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_SLATE

        y += Inches(0.95)

def slide_07_features_users(prs):
    """Platform Features for Users"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "geometric")
    add_title_bar(slide, "Platform Features", "For Users")

    # Retailers column
    ret_card = add_rounded_card(slide, Inches(0.5), Inches(1.7), Inches(4.4), Inches(5), PRIMARY_BLUE)
    ret_overlay = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(4.4), Inches(5))
    ret_overlay.fill.gradient()
    ret_overlay.fill.gradient_angle = 180
    ret_overlay.fill.gradient_stops[0].color.rgb = PRIMARY_BLUE
    ret_overlay.fill.gradient_stops[1].color.rgb = RGBColor(29, 78, 216)
    ret_overlay.line.fill.background()

    # Retailer icon
    ret_icon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(1.9), Inches(0.7), Inches(0.7))
    ret_icon.fill.solid()
    ret_icon.fill.fore_color.rgb = WHITE
    ret_icon.line.fill.background()

    ret_title = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(4), Inches(0.5))
    tf = ret_title.text_frame
    tf.text = "FOR RETAILERS"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    ret_features = [
        "Browse and compare products from multiple suppliers",
        "Place bulk orders with cart management",
        "Track order status in real-time",
        "Rate and review suppliers",
        "Receive personalized recommendations"
    ]

    y = Inches(3.4)
    for feat in ret_features:
        feat_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(3.8), Inches(0.6))
        tf = feat_box.text_frame
        tf.word_wrap = True
        tf.text = feat
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        y += Inches(0.65)

    # Suppliers column
    sup_card = add_rounded_card(slide, Inches(5.1), Inches(1.7), Inches(4.4), Inches(5), AMBER)
    sup_overlay = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.7), Inches(4.4), Inches(5))
    sup_overlay.fill.gradient()
    sup_overlay.fill.gradient_angle = 180
    sup_overlay.fill.gradient_stops[0].color.rgb = AMBER
    sup_overlay.fill.gradient_stops[1].color.rgb = ORANGE
    sup_overlay.line.fill.background()

    # Supplier icon
    sup_icon = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(6.8), Inches(1.9), Inches(0.7), Inches(0.7))
    sup_icon.fill.solid()
    sup_icon.fill.fore_color.rgb = WHITE
    sup_icon.line.fill.background()

    sup_title = slide.shapes.add_textbox(Inches(5.3), Inches(2.7), Inches(4), Inches(0.5))
    tf = sup_title.text_frame
    tf.text = "FOR SUPPLIERS"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sup_features = [
        "Manage product listings and inventory",
        "Approve or reject incoming orders",
        "Broadcast promotional announcements",
        "View demand analytics and reports"
    ]

    y = Inches(3.4)
    for feat in sup_features:
        feat_box = slide.shapes.add_textbox(Inches(5.4), y, Inches(3.8), Inches(0.7))
        tf = feat_box.text_frame
        tf.word_wrap = True
        tf.text = feat
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        y += Inches(0.75)

def slide_08_smart_features(prs):
    """Smart Platform Features"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "curves")
    add_title_bar(slide, "Smart Capabilities", "Machine Learning Powered")

    # ML Feature 1: Supplier Recommendation
    ml1_card = add_rounded_card(slide, Inches(0.6), Inches(1.7), Inches(8.8), Inches(1.4), PURPLE)
    ml1_overlay = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.7), Inches(8.8), Inches(1.4))
    ml1_overlay.fill.gradient()
    ml1_overlay.fill.gradient_angle = 90
    ml1_overlay.fill.gradient_stops[0].color.rgb = PURPLE
    ml1_overlay.fill.gradient_stops[1].color.rgb = RGBColor(109, 40, 217)
    ml1_overlay.line.fill.background()

    ml1_icon = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(1), Inches(1.95), Inches(0.9), Inches(0.9))
    ml1_icon.fill.solid()
    ml1_icon.fill.fore_color.rgb = WHITE
    ml1_icon.line.fill.background()

    ml1_title = slide.shapes.add_textbox(Inches(2.1), Inches(1.85), Inches(7), Inches(0.4))
    tf = ml1_title.text_frame
    tf.text = "Supplier Recommendation System"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    ml1_desc = slide.shapes.add_textbox(Inches(2.1), Inches(2.35), Inches(7), Inches(0.5))
    tf = ml1_desc.text_frame
    tf.text = "Ranks suppliers based on price, distance, performance, and preferences using Random Forest"
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(220, 200, 255)

    # ML Feature 2: Demand Forecasting
    ml2_card = add_rounded_card(slide, Inches(0.6), Inches(3.3), Inches(8.8), Inches(1.4), TEAL)
    ml2_overlay = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.3), Inches(8.8), Inches(1.4))
    ml2_overlay.fill.gradient()
    ml2_overlay.fill.gradient_angle = 90
    ml2_overlay.fill.gradient_stops[0].color.rgb = TEAL
    ml2_overlay.fill.gradient_stops[1].color.rgb = RGBColor(13, 148, 136)
    ml2_overlay.line.fill.background()

    ml2_icon = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(1), Inches(3.55), Inches(0.9), Inches(0.9))
    ml2_icon.fill.solid()
    ml2_icon.fill.fore_color.rgb = WHITE
    ml2_icon.line.fill.background()

    ml2_title = slide.shapes.add_textbox(Inches(2.1), Inches(3.45), Inches(7), Inches(0.4))
    tf = ml2_title.text_frame
    tf.text = "Demand Forecasting"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    ml2_desc = slide.shapes.add_textbox(Inches(2.1), Inches(3.95), Inches(7), Inches(0.5))
    tf = ml2_desc.text_frame
    tf.text = "Predicts future demand using historical data, moving averages, and seasonality patterns"
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(200, 240, 235)

    # Additional features grid
    add_features = [
        ("Real-time GPS Tracking", MSO_SHAPE.BEVEL, PRIMARY_BLUE),
        ("Chapa Payment Gateway", MSO_SHAPE.DIAMOND, ROSE),
        ("Automated Notifications", MSO_SHAPE.CHEVRON, GREEN),
        ("In-App Messaging", MSO_SHAPE.OVAL, AMBER),
    ]

    x_start = Inches(0.6)
    y = Inches(5)
    card_w = Inches(2.1)
    card_h = Inches(1.3)
    gap = Inches(0.15)

    for i, (name, icon_shape, color) in enumerate(add_features):
        x = x_start + i * (card_w + gap)
        card = add_rounded_card(slide, x, y, card_w, card_h, WHITE, color, 3)

        icon = slide.shapes.add_shape(icon_shape, x + Inches(0.7), y + Inches(0.2), Inches(0.6), Inches(0.6))
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.line.fill.background()

        name_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.85), card_w - Inches(0.2), Inches(0.4))
        tf = name_box.text_frame
        tf.word_wrap = True
        tf.text = name
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_SLATE

def slide_09_architecture(prs):
    """System Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_title_bar(slide, "Three-Tier Architecture")

    layers = [
        ("PRESENTATION LAYER", "Mobile App (Android)  •  Web Application  •  User Interface", PRIMARY_BLUE, Inches(1.7)),
        ("APPLICATION LAYER", "Node.js + Express  •  Authentication  •  Business Logic  •  ML Integration", TEAL, Inches(3.5)),
        ("DATA LAYER", "MySQL Database  •  User Data  •  Products & Orders  •  Analytics", AMBER, Inches(5.3)),
    ]

    for name, components, color, y_pos in layers:
        # Layer card
        layer_card = add_rounded_card(slide, Inches(1.2), y_pos, Inches(7.6), Inches(1.5), color)
        layer_overlay = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), y_pos, Inches(7.6), Inches(1.5))
        layer_overlay.fill.gradient()
        layer_overlay.fill.gradient_angle = 90
        layer_overlay.fill.gradient_stops[0].color.rgb = color
        layer_overlay.fill.gradient_stops[1].color.rgb = RGBColor(
            max(0, (color[0] if isinstance(color, tuple) else 30)),
            max(0, (color[1] if isinstance(color, tuple) else 100)),
            max(0, (color[2] if isinstance(color, tuple) else 180))
        )
        layer_overlay.line.fill.background()

        # Name
        name_box = slide.shapes.add_textbox(Inches(1.4), y_pos + Inches(0.2), Inches(7.2), Inches(0.5))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Components
        comp_box = slide.shapes.add_textbox(Inches(1.4), y_pos + Inches(0.75), Inches(7.2), Inches(0.6))
        tf = comp_box.text_frame
        tf.word_wrap = True
        tf.text = components
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(230, 240, 255)

    # Arrows between layers
    for y in [Inches(3.25), Inches(5.05)]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.5), y, Inches(1), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = SLATE
        arrow.line.fill.background()

def slide_10_tech_stack(prs):
    """Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "geometric")
    add_title_bar(slide, "Technology Stack")

    tech_data = [
        ("Frontend", "React.js with TypeScript  •  Tailwind CSS  •  Zustand", PRIMARY_BLUE, MSO_SHAPE.BEVEL),
        ("Backend", "Node.js with Express  •  JWT Authentication  •  RESTful APIs", TEAL, MSO_SHAPE.PENTAGON),
        ("Database", "MySQL with Sequelize ORM", AMBER, MSO_SHAPE.CAN),
        ("Machine Learning", "Python  •  Scikit-learn  •  Pandas  •  NumPy", PURPLE, MSO_SHAPE.PARALLELOGRAM),
        ("Payment", "Chapa Payment Gateway", ROSE, MSO_SHAPE.DIAMOND),
    ]

    y = Inches(1.7)
    for category, items, color, icon_shape in tech_data:
        # Card
        card = add_rounded_card(slide, Inches(0.8), y, Inches(8.4), Inches(0.9), WHITE, color, 3)

        # Icon
        icon = slide.shapes.add_shape(icon_shape, Inches(1.1), y + Inches(0.15), Inches(0.6), Inches(0.6))
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.line.fill.background()

        # Category
        cat_box = slide.shapes.add_textbox(Inches(1.9), y + Inches(0.15), Inches(2), Inches(0.6))
        tf = cat_box.text_frame
        tf.text = category
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color

        # Items
        items_box = slide.shapes.add_textbox(Inches(4), y + Inches(0.2), Inches(5), Inches(0.5))
        tf = items_box.text_frame
        tf.word_wrap = True
        tf.text = items
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = SLATE

        y += Inches(1.05)

def slide_11_user_roles(prs):
    """User Roles"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "curves")
    add_title_bar(slide, "Multi-Role Access Control")

    roles = [
        ("Retailer", "Browse products, Place orders, Track deliveries, Rate suppliers", PRIMARY_BLUE, MSO_SHAPE.RECTANGLE),
        ("Factory", "Manage products, Approve orders, View demand forecasts", TEAL, MSO_SHAPE.PENTAGON),
        ("Distributor", "Buy & sell, Manage inventory, Fulfill orders", GREEN, MSO_SHAPE.CUBE),
        ("Driver", "View assignments, Update delivery status, Track routes", AMBER, MSO_SHAPE.RIGHT_ARROW),
        ("Admin", "Manage users, Approve suppliers, Monitor platform", PURPLE, MSO_SHAPE.DIAMOND),
    ]

    y = Inches(1.7)
    for role, permissions, color, icon_shape in roles:
        # Card
        card = add_rounded_card(slide, Inches(0.8), y, Inches(8.4), Inches(0.9), WHITE, color, 3)

        # Icon
        icon = slide.shapes.add_shape(icon_shape, Inches(1.1), y + Inches(0.15), Inches(0.6), Inches(0.6))
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.line.fill.background()

        # Role
        role_box = slide.shapes.add_textbox(Inches(1.9), y + Inches(0.15), Inches(1.8), Inches(0.6))
        tf = role_box.text_frame
        tf.text = role
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color

        # Permissions
        perm_box = slide.shapes.add_textbox(Inches(3.8), y + Inches(0.2), Inches(5.2), Inches(0.5))
        tf = perm_box.text_frame
        tf.word_wrap = True
        tf.text = permissions
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = SLATE

        y += Inches(1)

def slide_12_database(prs):
    """Database Design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "circles")
    add_title_bar(slide, "Core Data Entities")

    entities = [
        ("Users", "Retailers, suppliers,\ndrivers, admins", MSO_SHAPE.OVAL),
        ("Products", "Name, price, stock,\nMOQ, images", MSO_SHAPE.CUBE),
        ("Orders", "Status, items,\ntracking, payments", MSO_SHAPE.BEVEL),
        ("Messages", "In-app chat\nhistory", MSO_SHAPE.CHEVRON),
        ("Ratings", "Supplier performance\nmetrics", MSO_SHAPE.DIAMOND),
        ("Payments", "Transactions,\nmethods, status", MSO_SHAPE.DIAMOND),
    ]

    colors = [PRIMARY_BLUE, TEAL, GREEN, PURPLE, AMBER, ROSE]

    x_positions = [Inches(0.6), Inches(3.5), Inches(6.4)]
    y_positions = [Inches(1.8), Inches(4.4)]

    for i, (name, desc, icon_shape) in enumerate(entities):
        row = i // 3
        col = i % 3
        x = x_positions[col]
        y = y_positions[row]
        color = colors[i]

        # Card
        card = add_rounded_card(slide, x, y, Inches(2.7), Inches(2.3), color)

        # Icon
        icon = slide.shapes.add_shape(icon_shape, x + Inches(0.9), y + Inches(0.3), Inches(0.9), Inches(0.9))
        icon.fill.solid()
        icon.fill.fore_color.rgb = WHITE
        icon.line.fill.background()

        # Name
        name_box = slide.shapes.add_textbox(x, y + Inches(1.3), Inches(2.7), Inches(0.4))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Description
        desc_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.7), Inches(2.5), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = desc
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(230, 240, 255)

def slide_13_ml_detail(prs):
    """ML Implementation Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_title_bar(slide, "Machine Learning Pipeline")

    # Process steps
    steps = [
        ("1", "DATA\nCOLLECTION", "Price, delivery time,\nratings, location", PRIMARY_BLUE),
        ("2", "MODEL\nTRAINING", "Random Forest\nClassifier", TEAL),
        ("3", "SCORE\nGENERATION", "Rank suppliers\nper retailer", GREEN),
        ("4", "PERSONALIZE", "User preference\nadaptation", PURPLE),
    ]

    x_start = Inches(0.6)
    y = Inches(1.8)
    box_w = Inches(2.1)
    box_h = Inches(2.3)
    gap = Inches(0.3)

    for i, (num, title, desc, color) in enumerate(steps):
        x = x_start + i * (box_w + gap)

        # Card
        card = add_rounded_card(slide, x, y, box_w, box_h, color)

        # Number
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), y + Inches(0.2), Inches(0.7), Inches(0.7))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = WHITE
        num_circle.line.fill.background()

        tf = num_circle.text_frame
        tf.text = num
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = color

        # Title
        title_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1), box_w - Inches(0.2), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Description
        desc_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.6), box_w - Inches(0.2), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = desc
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(220, 230, 255)

        # Arrow
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x + box_w + Inches(0.05), y + Inches(0.9), Inches(0.2), Inches(0.5))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = SLATE
            arrow.line.fill.background()

    # Features box
    feat_card = add_rounded_card(slide, Inches(0.6), Inches(4.4), Inches(8.8), Inches(2.2), WHITE, PRIMARY_BLUE, 3)

    feat_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(8.4), Inches(0.4))
    tf = feat_title.text_frame
    tf.text = "Features Used for Ranking"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    features = [
        "Price competitiveness",
        "On-time delivery rate",
        "Quality ratings",
        "Fulfillment time",
        "Communication responsiveness"
    ]

    x_positions = [Inches(1), Inches(5)]
    y_start = Inches(5.2)

    for i, feat in enumerate(features):
        col = i // 3
        row = i % 3
        x = x_positions[col]
        y = y_start + row * Inches(0.45)

        # Checkmark
        check = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.3), Inches(0.3))
        check.fill.solid()
        check.fill.fore_color.rgb = GREEN
        check.line.fill.background()

        # Text
        feat_box = slide.shapes.add_textbox(x + Inches(0.45), y, Inches(3.5), Inches(0.35))
        tf = feat_box.text_frame
        tf.text = feat
        p = tf.paragraphs[0]
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_SLATE

def slide_14_benefits(prs):
    """Benefits and Impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "geometric")
    add_title_bar(slide, "Expected Impact")

    columns = [
        ("For Retailers", [
            "Faster procurement process",
            "Better price comparison",
            "Improved supplier visibility",
            "Personalized recommendations"
        ], PRIMARY_BLUE),
        ("For Suppliers", [
            "Expanded market reach",
            "Direct buyer connections",
            "Reduced manual operations",
            "Access to demand insights"
        ], TEAL),
        ("For Industry", [
            "Digital transformation",
            "Reduced inefficiencies",
            "Increased transparency",
            "Data-driven decisions"
        ], AMBER),
    ]

    x_start = Inches(0.5)
    col_w = Inches(3)
    gap = Inches(0.15)

    for i, (title, items, color) in enumerate(columns):
        x = x_start + i * (col_w + gap)

        # Card
        card = add_rounded_card(slide, x, Inches(1.7), col_w, Inches(5), WHITE, color, 4)

        # Header
        header = add_rounded_card(slide, x, Inches(1.7), col_w, Inches(0.7), color)
        header_text = slide.shapes.add_textbox(x, Inches(1.85), col_w, Inches(0.4))
        tf = header_text.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Items
        y = Inches(2.7)
        for item in items:
            # Bullet
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.1), Inches(0.25), Inches(0.25))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = color
            bullet.line.fill.background()

            # Text
            item_box = slide.shapes.add_textbox(x + Inches(0.7), y, col_w - Inches(0.9), Inches(0.6))
            tf = item_box.text_frame
            tf.word_wrap = True
            tf.text = item
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_SLATE

            y += Inches(0.75)

def slide_15_feasibility(prs):
    """Feasibility Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "curves")
    add_title_bar(slide, "Project Feasibility")

    feasibility_data = [
        ("Technical", [
            "Proven technologies (React, Node.js, MySQL)",
            "Team has required skills",
            "No specialized hardware needed"
        ], PRIMARY_BLUE, MSO_SHAPE.PENTAGON),
        ("Operational", [
            "User-friendly interface design",
            "Supports existing business practices",
            "Mobile-first approach"
        ], TEAL, MSO_SHAPE.OVAL),
        ("Economic", [
            "Low development cost (~9,500 ETB)",
            "Open-source technologies",
            "Revenue through transaction fees"
        ], GREEN, MSO_SHAPE.DIAMOND),
    ]

    x_start = Inches(0.5)
    col_w = Inches(3)
    gap = Inches(0.15)

    for i, (title, items, color, icon_shape) in enumerate(feasibility_data):
        x = x_start + i * (col_w + gap)

        # Card
        card = add_rounded_card(slide, x, Inches(1.7), col_w, Inches(5), WHITE, color, 4)

        # Icon
        icon = slide.shapes.add_shape(icon_shape, x + Inches(1.1), Inches(1.9), Inches(0.8), Inches(0.8))
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(x, Inches(2.85), col_w, Inches(0.5))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color

        # Items
        y = Inches(3.5)
        for item in items:
            item_box = slide.shapes.add_textbox(x + Inches(0.2), y, col_w - Inches(0.4), Inches(0.8))
            tf = item_box.text_frame
            tf.word_wrap = True
            tf.text = f"• {item}"
            p = tf.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_SLATE
            y += Inches(0.7)

def slide_16_methodology(prs):
    """Development Methodology"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_title_bar(slide, "Development Approach", "Agile & Incremental Methodology")

    phases = [
        ("Planning", "Week 1-2", PRIMARY_BLUE),
        ("Analysis", "Week 3-8", TEAL),
        ("Design", "Week 8-12", GREEN),
        ("Implementation", "Week 13-15", PURPLE),
        ("Testing", "Week 15-16", AMBER),
        ("Documentation", "Week 17-18", ROSE),
    ]

    x_positions = [Inches(0.6), Inches(3.5), Inches(6.4)]
    y_positions = [Inches(1.8), Inches(4)]

    for i, (phase, timeline, color) in enumerate(phases):
        row = i // 3
        col = i % 3
        x = x_positions[col]
        y = y_positions[row]

        # Card
        card = add_rounded_card(slide, x, y, Inches(2.7), Inches(1.8), color)

        # Phase name
        phase_box = slide.shapes.add_textbox(x, y + Inches(0.5), Inches(2.7), Inches(0.5))
        tf = phase_box.text_frame
        tf.text = phase
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Timeline
        time_box = slide.shapes.add_textbox(x, y + Inches(1.1), Inches(2.7), Inches(0.4))
        tf = time_box.text_frame
        tf.text = timeline
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(220, 230, 255)

    # Timeline banner
    banner = add_rounded_card(slide, Inches(2.5), Inches(6.2), Inches(5), Inches(0.7), GREEN)
    banner_text = slide.shapes.add_textbox(Inches(2.5), Inches(6.35), Inches(5), Inches(0.4))
    tf = banner_text.text_frame
    tf.text = "Total Timeline: ~18 weeks (2 semesters)"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE

def slide_17_challenges(prs):
    """Challenges and Mitigations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "circles")
    add_title_bar(slide, "Challenges & Mitigations")

    challenges = [
        "Limited initial ML training data",
        "Dependency on stable internet connectivity",
        "Resistance to digital adoption",
        "Reliance on third-party payment providers",
    ]

    y = Inches(1.7)
    for challenge in challenges:
        # Card
        card = add_rounded_card(slide, Inches(0.8), y, Inches(8.4), Inches(0.65), WHITE, ROSE, 2)

        # X icon
        x_icon = slide.shapes.add_shape(MSO_SHAPE.CROSS, Inches(1.1), y + Inches(0.12), Inches(0.4), Inches(0.4))
        x_icon.fill.solid()
        x_icon.fill.fore_color.rgb = ROSE
        x_icon.line.fill.background()

        # Text
        text_box = slide.shapes.add_textbox(Inches(1.7), y + Inches(0.12), Inches(7.3), Inches(0.4))
        tf = text_box.text_frame
        tf.text = challenge
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_SLATE

        y += Inches(0.8)

    # Mitigation box
    mit_card = add_rounded_card(slide, Inches(0.8), Inches(5), Inches(8.4), Inches(1.7), TEAL)

    mit_title = slide.shapes.add_textbox(Inches(1), Inches(5.15), Inches(8), Inches(0.4))
    tf = mit_title.text_frame
    tf.text = "MITIGATION STRATEGIES"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE

    mitigations = [
        "Hybrid ML approach (rule-based + data-driven)",
        "Offline capabilities for critical functions",
        "Comprehensive user onboarding and support"
    ]

    y = Inches(5.7)
    for mit in mitigations:
        mit_box = slide.shapes.add_textbox(Inches(1.3), y, Inches(7.5), Inches(0.35))
        tf = mit_box.text_frame
        tf.text = f"   {mit}"
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        y += Inches(0.4)

def slide_18_future(prs):
    """Future Enhancements"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    add_decorative_background(slide, "geometric")
    add_title_bar(slide, "Future Roadmap")

    enhancements = [
        ("Additional product categories", MSO_SHAPE.CUBE),
        ("AI-powered chatbot support", MSO_SHAPE.CHEVRON),
        ("iOS mobile application", MSO_SHAPE.BEVEL),
        ("Multi-language support", MSO_SHAPE.CLOUD),
        ("Advanced analytics dashboards", MSO_SHAPE.PARALLELOGRAM),
        ("ERP system integration", MSO_SHAPE.PENTAGON),
        ("Logistics optimization", MSO_SHAPE.RIGHT_ARROW),
        ("Blockchain transparency", MSO_SHAPE.HEXAGON),
    ]

    colors = [PRIMARY_BLUE, TEAL, GREEN, PURPLE, AMBER, ROSE, PRIMARY_BLUE, TEAL]

    x_positions = [Inches(0.6), Inches(5)]
    y_start = Inches(1.8)
    card_w = Inches(4.4)
    card_h = Inches(0.85)

    for i, (text, icon_shape) in enumerate(enhancements):
        row = i // 2
        col = i % 2
        x = x_positions[col]
        y = y_start + row * (card_h + Inches(0.2))
        color = colors[i]

        # Card
        card = add_rounded_card(slide, x, y, card_w, card_h, WHITE, color, 2)

        # Icon
        icon = slide.shapes.add_shape(icon_shape, x + Inches(0.2), y + Inches(0.15), Inches(0.55), Inches(0.55))
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.line.fill.background()

        # Text
        text_box = slide.shapes.add_textbox(x + Inches(0.9), y + Inches(0.2), card_w - Inches(1.1), Inches(0.45))
        tf = text_box.text_frame
        tf.text = text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_SLATE

def slide_19_conclusion(prs):
    """Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background
    bg.fill.gradient()
    bg.fill.gradient_angle = 135
    bg.fill.gradient_stops[0].color.rgb = DEEP_BLUE
    bg.fill.gradient_stops[1].color.rgb = PRIMARY_BLUE

    # Decorative
    wave = slide.shapes.add_shape(MSO_SHAPE.WAVE, Inches(6), Inches(4.5), Inches(5), Inches(3.5))
    wave.fill.solid()
    wave.fill.fore_color.rgb = RGBColor(30, 64, 175)
    wave.line.fill.background()

    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    tf = title.text_frame
    tf.text = "Summary"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Deliverables card
    deliv_card = add_rounded_card(slide, Inches(1.2), Inches(1.8), Inches(7.6), Inches(3.5), WHITE)

    deliv_title = slide.shapes.add_textbox(Inches(1.4), Inches(2), Inches(7.2), Inches(0.5))
    tf = deliv_title.text_frame
    tf.text = "TradeBridge Delivers:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    deliverables = [
        "Centralized B2B marketplace for Ethiopian wholesale",
        "Smart supplier recommendations & demand forecasting",
        "Real-time tracking & secure payments",
        "Improved efficiency across the supply chain"
    ]

    y = Inches(2.7)
    for item in deliverables:
        # Check
        check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.6), y + Inches(0.05), Inches(0.3), Inches(0.3))
        check.fill.solid()
        check.fill.fore_color.rgb = GREEN
        check.line.fill.background()

        # Text
        item_box = slide.shapes.add_textbox(Inches(2.1), y, Inches(6.5), Inches(0.4))
        tf = item_box.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_SLATE
        y += Inches(0.5)

    # Status banner
    status = add_rounded_card(slide, Inches(2), Inches(5.6), Inches(6), Inches(0.7), GREEN)
    status_text = slide.shapes.add_textbox(Inches(2), Inches(5.75), Inches(6), Inches(0.4))
    tf = status_text.text_frame
    tf.text = "Prototype Completed - Ready for Deployment"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE

def slide_20_thankyou(prs):
    """Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background
    bg.fill.gradient()
    bg.fill.gradient_angle = 135
    bg.fill.gradient_stops[0].color.rgb = TEAL
    bg.fill.gradient_stops[1].color.rgb = RGBColor(13, 148, 136)

    # Decorative circles
    add_circle_accent(slide, Inches(8), Inches(0.3), Inches(1.5), RGBColor(45, 212, 191))
    add_circle_accent(slide, Inches(-0.5), Inches(5), Inches(2), RGBColor(45, 212, 191))
    add_circle_accent(slide, Inches(9), Inches(5.5), Inches(1), RGBColor(45, 212, 191))

    # Thank you text
    thanks = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.2))
    tf = thanks.text_frame
    tf.text = "Thank You"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.5), Inches(3), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = AMBER
    line.line.fill.background()

    # Questions
    questions = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    tf = questions.text_frame
    tf.text = "Questions?"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE

    # Contact
    contact = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(6), Inches(0.8))
    tf = contact.text_frame
    tf.text = "Advisor: Dr. Ejigu Tefere\nAdama Science and Technology University"
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.color.rgb = WHITE

    # Team
    team = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.4))
    tf = team.text_frame
    tf.text = "Hidaya Nurmeka  |  Ebisa Gutema  |  Hana Kebede  |  Hana Jote  |  Ilham Mohammedhassen"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 240, 235)

def create_all_slides(prs):
    """Create all slides"""
    slides = [
        ("Title", slide_01_title),
        ("Agenda", slide_02_agenda),
        ("Problem Statement", slide_03_problem),
        ("Introduction", slide_04_intro),
        ("Scope", slide_05_scope),
        ("Objectives", slide_06_objectives),
        ("Features - Users", slide_07_features_users),
        ("Smart Features", slide_08_smart_features),
        ("Architecture", slide_09_architecture),
        ("Technology Stack", slide_10_tech_stack),
        ("User Roles", slide_11_user_roles),
        ("Database Design", slide_12_database),
        ("ML Pipeline", slide_13_ml_detail),
        ("Benefits", slide_14_benefits),
        ("Feasibility", slide_15_feasibility),
        ("Methodology", slide_16_methodology),
        ("Challenges", slide_17_challenges),
        ("Future Roadmap", slide_18_future),
        ("Conclusion", slide_19_conclusion),
        ("Thank You", slide_20_thankyou),
    ]

    for name, func in slides:
        print(f"  Creating: {name}")
        func(prs)

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("\nCreating TradeBridge Stunning Presentation...\n")
    create_all_slides(prs)

    output_path = 'TradeBridge_Stunning.pptx'
    prs.save(output_path)

    print(f"\nPresentation created successfully!")
    print(f"Location: {os.path.abspath(output_path)}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"File size: {os.path.getsize(output_path) / (1024 * 1024):.2f} MB\n")

if __name__ == '__main__':
    main()
