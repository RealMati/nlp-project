#!/usr/bin/env python3
"""
TradeBridge PowerPoint Presentation Generator
Creates a luxurious, visually appealing presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color Scheme
PRIMARY_BLUE = RGBColor(37, 99, 235)      # #2563EB
SECONDARY_GREEN = RGBColor(16, 185, 129)  # #10B981
ACCENT_ORANGE = RGBColor(245, 158, 11)    # #F59E0B
DARK_TEXT = RGBColor(31, 41, 55)          # #1F2937
LIGHT_GRAY = RGBColor(249, 250, 251)      # #F9FAFB
WHITE = RGBColor(255, 255, 255)

def add_gradient_background(slide, color1, color2):
    """Add a gradient background to slide"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45.0
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_gradient_background(slide, PRIMARY_BLUE, RGBColor(59, 130, 246))

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "TradeBridge"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(66)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Smart Supply–Demand Management System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = WHITE

    # Institution
    inst_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.8))
    inst_frame = inst_box.text_frame
    inst_frame.text = "Adama Science and Technology University\nCollege of Electrical Engineering and Computing"
    for para in inst_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.color.rgb = WHITE

    # Team members at bottom
    team_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
    team_frame = team_box.text_frame
    team_frame.text = "Hidaya Nurmeka • Ebisa Gutema • Hana Kebede • Hana Jote • Ilham Mohammedhassen"
    team_para = team_frame.paragraphs[0]
    team_para.alignment = PP_ALIGN.CENTER
    team_para.font.size = Pt(12)
    team_para.font.color.rgb = WHITE

def add_agenda_slide(prs):
    """Slide 2: Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Presentation Outline"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE

    # Agenda items
    agenda_items = [
        "Problem Overview",
        "System Objectives",
        "Key Features",
        "System Architecture",
        "Technology Stack",
        "Implementation Highlights",
        "Machine Learning Integration",
        "Benefits & Impact",
        "Demo/Prototype",
        "Conclusion"
    ]

    left_col = Inches(1)
    top_start = Inches(1.8)

    for i, item in enumerate(agenda_items):
        row = i % 5
        col = i // 5

        x_pos = left_col + (col * Inches(4.5))
        y_pos = top_start + (row * Inches(0.9))

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x_pos, y_pos,
            Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_ORANGE
        circle.line.color.rgb = ACCENT_ORANGE

        num_frame = circle.text_frame
        num_frame.text = str(i + 1)
        num_para = num_frame.paragraphs[0]
        num_para.alignment = PP_ALIGN.CENTER
        num_para.font.size = Pt(18)
        num_para.font.bold = True
        num_para.font.color.rgb = WHITE

        # Item text
        text_box = slide.shapes.add_textbox(x_pos + Inches(0.7), y_pos, Inches(3.5), Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = item
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(18)
        text_para.font.color.rgb = DARK_TEXT

def add_problem_slide(prs):
    """Slide 3: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Title with accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.5),
        Inches(0.3), Inches(0.8)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_ORANGE
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Current Supply Chain Challenges in Ethiopia"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE

    # Problems
    problems = [
        "Manual, time-consuming procurement processes",
        "Lack of centralized B2B marketplace",
        "Poor demand planning and forecasting",
        "Limited supplier visibility and comparison",
        "Inefficient communication between stakeholders",
        "Stock shortages and delivery delays"
    ]

    y_pos = Inches(2)
    for problem in problems:
        # X icon
        x_mark = slide.shapes.add_textbox(Inches(1), y_pos, Inches(0.5), Inches(0.5))
        x_frame = x_mark.text_frame
        x_frame.text = "❌"
        x_para = x_frame.paragraphs[0]
        x_para.font.size = Pt(24)

        # Problem text
        prob_box = slide.shapes.add_textbox(Inches(1.8), y_pos, Inches(7.5), Inches(0.5))
        prob_frame = prob_box.text_frame
        prob_frame.text = problem
        prob_para = prob_frame.paragraphs[0]
        prob_para.font.size = Pt(20)
        prob_para.font.color.rgb = DARK_TEXT

        y_pos += Inches(0.7)

def add_intro_slide(prs):
    """Slide 4: What is TradeBridge?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, WHITE, LIGHT_GRAY)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Introducing TradeBridge"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE

    # Description box with border
    desc_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.8),
        Inches(8), Inches(1.5)
    )
    desc_shape.fill.solid()
    desc_shape.fill.fore_color.rgb = WHITE
    desc_shape.line.color.rgb = PRIMARY_BLUE
    desc_shape.line.width = Pt(3)

    desc_frame = desc_shape.text_frame
    desc_frame.word_wrap = True
    desc_frame.text = "A comprehensive B2B digital platform that connects retailers, factories, distributors, and delivery personnel to streamline bulk ordering, improve supply chain transparency, and enable data-driven decision-making in the Ethiopian wholesale market."
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(20)
    desc_para.font.color.rgb = DARK_TEXT
    desc_para.alignment = PP_ALIGN.CENTER

    # Stakeholder boxes
    stakeholders = [
        ("🏪", "Retailers"),
        ("🏭", "Factories"),
        ("🚚", "Distributors"),
        ("📦", "Delivery\nPersonnel")
    ]

    x_start = Inches(1.5)
    y_pos = Inches(4)
    width = Inches(1.5)

    for i, (emoji, name) in enumerate(stakeholders):
        x_pos = x_start + (i * Inches(2))

        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos,
            width, Inches(1.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = SECONDARY_GREEN
        box.line.fill.background()

        # Emoji
        emoji_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.2), width, Inches(0.6))
        emoji_frame = emoji_box.text_frame
        emoji_frame.text = emoji
        emoji_para = emoji_frame.paragraphs[0]
        emoji_para.alignment = PP_ALIGN.CENTER
        emoji_para.font.size = Pt(36)

        # Name
        name_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.8), width, Inches(0.6))
        name_frame = name_box.text_frame
        name_frame.text = name
        name_para = name_frame.paragraphs[0]
        name_para.alignment = PP_ALIGN.CENTER
        name_para.font.size = Pt(16)
        name_para.font.bold = True
        name_para.font.color.rgb = WHITE

def add_scope_slide(prs):
    """Slide 5: Project Scope"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Scope & Focus"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE

    # In Scope box
    in_scope_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.8),
        Inches(4), Inches(4)
    )
    in_scope_box.fill.solid()
    in_scope_box.fill.fore_color.rgb = SECONDARY_GREEN
    in_scope_box.line.fill.background()

    in_title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3.6), Inches(0.5))
    in_title_frame = in_title.text_frame
    in_title_frame.text = "✅ In Scope"
    in_title_para = in_title_frame.paragraphs[0]
    in_title_para.alignment = PP_ALIGN.CENTER
    in_title_para.font.size = Pt(26)
    in_title_para.font.bold = True
    in_title_para.font.color.rgb = WHITE

    in_items = [
        "Food and beverage products",
        "Micro to large enterprises",
        "Ethiopian market (ETB)",
        "Mobile & Web platforms",
        "Distribution of finished goods"
    ]

    y_pos = Inches(2.8)
    for item in in_items:
        item_box = slide.shapes.add_textbox(Inches(1.2), y_pos, Inches(3.2), Inches(0.5))
        item_frame = item_box.text_frame
        item_frame.text = f"• {item}"
        item_para = item_frame.paragraphs[0]
        item_para.font.size = Pt(16)
        item_para.font.color.rgb = WHITE
        y_pos += Inches(0.6)

    # Out of Scope box
    out_scope_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.8),
        Inches(4), Inches(4)
    )
    out_scope_box.fill.solid()
    out_scope_box.fill.fore_color.rgb = RGBColor(239, 68, 68)  # Red
    out_scope_box.line.fill.background()

    out_title = slide.shapes.add_textbox(Inches(5.4), Inches(2), Inches(3.6), Inches(0.5))
    out_title_frame = out_title.text_frame
    out_title_frame.text = "❌ Out of Scope"
    out_title_para = out_title_frame.paragraphs[0]
    out_title_para.alignment = PP_ALIGN.CENTER
    out_title_para.font.size = Pt(26)
    out_title_para.font.bold = True
    out_title_para.font.color.rgb = WHITE

    out_items = [
        "Raw material procurement",
        "International trade",
        "Very large national producers"
    ]

    y_pos = Inches(2.8)
    for item in out_items:
        item_box = slide.shapes.add_textbox(Inches(5.6), y_pos, Inches(3.2), Inches(0.5))
        item_frame = item_box.text_frame
        item_frame.text = f"• {item}"
        item_para = item_frame.paragraphs[0]
        item_para.font.size = Pt(16)
        item_para.font.color.rgb = WHITE
        y_pos += Inches(0.6)

def add_objectives_slide(prs):
    """Slide 6: System Objectives"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Key Objectives"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE

    # General objective box
    gen_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(8), Inches(1)
    )
    gen_box.fill.solid()
    gen_box.fill.fore_color.rgb = PRIMARY_BLUE
    gen_box.line.fill.background()

    gen_frame = gen_box.text_frame
    gen_frame.word_wrap = True
    gen_frame.text = "General Objective: Design and develop a digital platform to streamline B2B procurement and enhance supply chain visibility."
    gen_para = gen_frame.paragraphs[0]
    gen_para.font.size = Pt(18)
    gen_para.font.color.rgb = WHITE
    gen_para.alignment = PP_ALIGN.CENTER

    # Specific objectives
    objectives = [
        ("📱", "Develop Web & Mobile Application"),
        ("🤖", "Implement ML-based Supplier Recommendation"),
        ("📊", "Introduce Demand Forecasting Capabilities"),
        ("💬", "Enable Real-time Communication"),
        ("🔍", "Provide Centralized Supplier Directory")
    ]

    y_start = Inches(3)
    for i, (emoji, obj) in enumerate(objectives):
        y_pos = y_start + (i * Inches(0.85))

        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(1.5), y_pos, Inches(0.6), Inches(0.5))
        emoji_frame = emoji_box.text_frame
        emoji_frame.text = emoji
        emoji_para = emoji_frame.paragraphs[0]
        emoji_para.font.size = Pt(28)

        # Objective text
        obj_box = slide.shapes.add_textbox(Inches(2.3), y_pos, Inches(6.5), Inches(0.6))
        obj_frame = obj_box.text_frame
        obj_frame.text = obj
        obj_para = obj_frame.paragraphs[0]
        obj_para.font.size = Pt(22)
        obj_para.font.color.rgb = DARK_TEXT

def add_architecture_slide(prs):
    """Slide 9: System Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, WHITE, LIGHT_GRAY)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Three-Tier Architecture"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE

    layers = [
        ("Presentation Layer", ["Mobile App (Android)", "Web Application", "User Interface Components"], PRIMARY_BLUE),
        ("Application Layer", ["Authentication & Authorization", "Business Logic", "API Services", "ML Model Integration"], SECONDARY_GREEN),
        ("Data Layer", ["User Data", "Products & Orders", "Analytics & Logs"], ACCENT_ORANGE)
    ]

    y_pos = Inches(1.8)
    for layer_name, components, color in layers:
        # Layer box
        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), y_pos,
            Inches(7), Inches(1.5)
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = color
        layer_box.line.fill.background()

        # Layer name
        name_box = slide.shapes.add_textbox(Inches(1.7), y_pos + Inches(0.1), Inches(6.6), Inches(0.4))
        name_frame = name_box.text_frame
        name_frame.text = layer_name
        name_para = name_frame.paragraphs[0]
        name_para.alignment = PP_ALIGN.CENTER
        name_para.font.size = Pt(24)
        name_para.font.bold = True
        name_para.font.color.rgb = WHITE

        # Components
        comp_text = " • ".join(components)
        comp_box = slide.shapes.add_textbox(Inches(1.7), y_pos + Inches(0.6), Inches(6.6), Inches(0.8))
        comp_frame = comp_box.text_frame
        comp_frame.word_wrap = True
        comp_frame.text = comp_text
        comp_para = comp_frame.paragraphs[0]
        comp_para.alignment = PP_ALIGN.CENTER
        comp_para.font.size = Pt(14)
        comp_para.font.color.rgb = WHITE

        # Arrow
        if y_pos < Inches(5):
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(4.5), y_pos + Inches(1.6),
                Inches(1), Inches(0.6)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_TEXT
            arrow.line.fill.background()

        y_pos += Inches(2.2)

def add_tech_stack_slide(prs):
    """Slide 10: Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Technologies Used"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE

    tech_categories = [
        ("Frontend", ["⚛️ React.js with TypeScript", "🎨 Tailwind CSS", "📊 Zustand (State Management)"], PRIMARY_BLUE),
        ("Backend", ["🟢 Node.js with Express", "🔐 JWT Authentication", "📡 RESTful APIs"], SECONDARY_GREEN),
        ("Database", ["🗄️ MySQL with Sequelize ORM"], ACCENT_ORANGE),
        ("Machine Learning", ["🐍 Python", "📚 Scikit-learn, Pandas, NumPy"], RGBColor(139, 92, 246)),
        ("Payment", ["💰 Chapa Payment Gateway"], RGBColor(236, 72, 153))
    ]

    y_pos = Inches(1.8)
    for category, items, color in tech_categories:
        # Category box
        cat_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), y_pos,
            Inches(8), Inches(0.9)
        )
        cat_box.fill.solid()
        cat_box.fill.fore_color.rgb = color
        cat_box.line.fill.background()

        # Category name
        cat_name = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.1), Inches(2), Inches(0.3))
        cat_frame = cat_name.text_frame
        cat_frame.text = category
        cat_para = cat_frame.paragraphs[0]
        cat_para.font.size = Pt(20)
        cat_para.font.bold = True
        cat_para.font.color.rgb = WHITE

        # Items
        items_text = " • ".join(items)
        items_box = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.45), Inches(7.6), Inches(0.4))
        items_frame = items_box.text_frame
        items_frame.text = items_text
        items_para = items_frame.paragraphs[0]
        items_para.font.size = Pt(16)
        items_para.font.color.rgb = WHITE

        y_pos += Inches(1.05)

def add_ml_recommendation_slide(prs):
    """Slide 13: ML Supplier Recommendation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Intelligent Supplier Ranking"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE

    # Process flow
    steps = [
        ("1", "Collects Data", "Price, delivery time, ratings, location, order history", PRIMARY_BLUE),
        ("2", "Trains Model", "Random Forest Classifier", SECONDARY_GREEN),
        ("3", "Generates Score", "Ranks suppliers for each retailer", ACCENT_ORANGE),
        ("4", "Personalizes", "Based on retailer's past preferences", RGBColor(139, 92, 246))
    ]

    x_start = Inches(0.8)
    y_pos = Inches(2)
    box_width = Inches(2)

    for i, (num, title, desc, color) in enumerate(steps):
        x_pos = x_start + (i * Inches(2.3))

        # Step box
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos,
            box_width, Inches(1.5)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = color
        step_box.line.fill.background()

        # Number
        num_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.15), box_width, Inches(0.4))
        num_frame = num_box.text_frame
        num_frame.text = num
        num_para = num_frame.paragraphs[0]
        num_para.alignment = PP_ALIGN.CENTER
        num_para.font.size = Pt(32)
        num_para.font.bold = True
        num_para.font.color.rgb = WHITE

        # Title
        title_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.6), box_width, Inches(0.3))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE

        # Description
        desc_box = slide.shapes.add_textbox(x_pos + Inches(0.1), y_pos + Inches(0.95), box_width - Inches(0.2), Inches(0.5))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_frame.text = desc
        desc_para = desc_frame.paragraphs[0]
        desc_para.alignment = PP_ALIGN.CENTER
        desc_para.font.size = Pt(11)
        desc_para.font.color.rgb = WHITE

        # Arrow
        if i < 3:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x_pos + box_width + Inches(0.05), y_pos + Inches(0.6),
                Inches(0.2), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_TEXT
            arrow.line.fill.background()

    # Features used section
    features_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(4),
        Inches(8), Inches(2.3)
    )
    features_box.fill.solid()
    features_box.fill.fore_color.rgb = LIGHT_GRAY
    features_box.line.color.rgb = PRIMARY_BLUE
    features_box.line.width = Pt(2)

    feat_title = slide.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(7.6), Inches(0.4))
    feat_title_frame = feat_title.text_frame
    feat_title_frame.text = "Features Used for Ranking:"
    feat_title_para = feat_title_frame.paragraphs[0]
    feat_title_para.font.size = Pt(20)
    feat_title_para.font.bold = True
    feat_title_para.font.color.rgb = PRIMARY_BLUE

    features = [
        "✓ Price competitiveness",
        "✓ On-time delivery rate",
        "✓ Quality ratings",
        "✓ Fulfillment time",
        "✓ Communication responsiveness"
    ]

    y_feat = Inches(4.8)
    for feature in features:
        feat_box = slide.shapes.add_textbox(Inches(1.5), y_feat, Inches(7), Inches(0.3))
        feat_frame = feat_box.text_frame
        feat_frame.text = feature
        feat_para = feat_frame.paragraphs[0]
        feat_para.font.size = Pt(18)
        feat_para.font.color.rgb = DARK_TEXT
        y_feat += Inches(0.35)

def add_benefits_slide(prs):
    """Slide 17: Benefits & Impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, LIGHT_GRAY, WHITE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Expected Impact"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE

    # Three columns
    columns = [
        ("For Retailers", [
            "⚡ Faster procurement process",
            "💰 Better price comparison",
            "📊 Improved supplier visibility",
            "🎯 Personalized recommendations"
        ], PRIMARY_BLUE),
        ("For Suppliers", [
            "📈 Expanded market reach",
            "🤝 Direct buyer connections",
            "📉 Reduced manual operations",
            "📊 Access to demand insights"
        ], SECONDARY_GREEN),
        ("For the Industry", [
            "🌐 Digital transformation",
            "📉 Reduced inefficiencies",
            "🔍 Increased transparency",
            "📈 Data-driven decisions"
        ], ACCENT_ORANGE)
    ]

    x_start = Inches(0.5)
    col_width = Inches(3)

    for i, (title, benefits, color) in enumerate(columns):
        x_pos = x_start + (i * Inches(3.2))

        # Column header
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, Inches(1.5),
            col_width, Inches(0.6)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = color
        header_box.line.fill.background()

        header_frame = header_box.text_frame
        header_frame.text = title
        header_para = header_frame.paragraphs[0]
        header_para.alignment = PP_ALIGN.CENTER
        header_para.font.size = Pt(22)
        header_para.font.bold = True
        header_para.font.color.rgb = WHITE

        # Benefits
        y_pos = Inches(2.3)
        for benefit in benefits:
            ben_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y_pos, col_width - Inches(0.4), Inches(0.5))
            ben_frame = ben_box.text_frame
            ben_frame.word_wrap = True
            ben_frame.text = benefit
            ben_para = ben_frame.paragraphs[0]
            ben_para.font.size = Pt(16)
            ben_para.font.color.rgb = DARK_TEXT
            y_pos += Inches(0.75)

def add_conclusion_slide(prs):
    """Slide 22: Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, PRIMARY_BLUE, RGBColor(59, 130, 246))

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Summary"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Deliverables box
    deliver_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(2),
        Inches(7), Inches(2)
    )
    deliver_box.fill.solid()
    deliver_box.fill.fore_color.rgb = WHITE
    deliver_box.line.fill.background()

    deliver_title = slide.shapes.add_textbox(Inches(1.7), Inches(2.2), Inches(6.6), Inches(0.4))
    deliver_title_frame = deliver_title.text_frame
    deliver_title_frame.text = "TradeBridge delivers:"
    deliver_title_para = deliver_title_frame.paragraphs[0]
    deliver_title_para.font.size = Pt(24)
    deliver_title_para.font.bold = True
    deliver_title_para.font.color.rgb = PRIMARY_BLUE

    deliverables = [
        "✅ Centralized B2B marketplace for Ethiopian wholesale",
        "✅ Smart supplier recommendations & demand forecasting",
        "✅ Real-time tracking & secure payments",
        "✅ Improved efficiency across the supply chain"
    ]

    y_pos = Inches(2.8)
    for item in deliverables:
        item_box = slide.shapes.add_textbox(Inches(2), y_pos, Inches(6), Inches(0.3))
        item_frame = item_box.text_frame
        item_frame.text = item
        item_para = item_frame.paragraphs[0]
        item_para.font.size = Pt(16)
        item_para.font.color.rgb = DARK_TEXT
        y_pos += Inches(0.4)

    # Status banner
    status_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(4.8),
        Inches(6), Inches(0.7)
    )
    status_box.fill.solid()
    status_box.fill.fore_color.rgb = SECONDARY_GREEN
    status_box.line.fill.background()

    status_frame = status_box.text_frame
    status_frame.text = "Status: Prototype completed, ready for deployment"
    status_para = status_frame.paragraphs[0]
    status_para.alignment = PP_ALIGN.CENTER
    status_para.font.size = Pt(20)
    status_para.font.bold = True
    status_para.font.color.rgb = WHITE

def add_thank_you_slide(prs):
    """Slide 23: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, SECONDARY_GREEN, RGBColor(5, 150, 105))

    # Thank you
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You!"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.alignment = PP_ALIGN.CENTER
    thanks_para.font.size = Pt(72)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = WHITE

    # Questions
    q_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    q_frame = q_box.text_frame
    q_frame.text = "Questions?"
    q_para = q_frame.paragraphs[0]
    q_para.alignment = PP_ALIGN.CENTER
    q_para.font.size = Pt(36)
    q_para.font.color.rgb = WHITE

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Advisor: Dr. Ejigu Tefere\nAdama Science and Technology University"
    for para in contact_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = WHITE

def create_presentation():
    """Main function to create the presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating TradeBridge Presentation...")

    print("  Adding Slide 1: Title")
    add_title_slide(prs)

    print("  Adding Slide 2: Agenda")
    add_agenda_slide(prs)

    print("  Adding Slide 3: Problem Statement")
    add_problem_slide(prs)

    print("  Adding Slide 4: Introduction")
    add_intro_slide(prs)

    print("  Adding Slide 5: Scope")
    add_scope_slide(prs)

    print("  Adding Slide 6: Objectives")
    add_objectives_slide(prs)

    print("  Adding Slide 9: Architecture")
    add_architecture_slide(prs)

    print("  Adding Slide 10: Technology Stack")
    add_tech_stack_slide(prs)

    print("  Adding Slide 13: ML Recommendation")
    add_ml_recommendation_slide(prs)

    print("  Adding Slide 17: Benefits")
    add_benefits_slide(prs)

    print("  Adding Slide 22: Conclusion")
    add_conclusion_slide(prs)

    print("  Adding Slide 23: Thank You")
    add_thank_you_slide(prs)

    # Save
    output_path = '/Users/leul/Documents/NLP/proj/TradeBridge_Presentation.pptx'
    prs.save(output_path)
    print(f"\n✅ Presentation created successfully!")
    print(f"📁 Location: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == '__main__':
    create_presentation()
