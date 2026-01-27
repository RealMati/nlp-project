#!/usr/bin/env python3
"""
TradeBridge Premium PowerPoint Presentation Generator
Creates a comprehensive, luxurious presentation with all 23 slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

# Premium Color Palette
PRIMARY_BLUE = RGBColor(37, 99, 235)
BLUE_LIGHT = RGBColor(59, 130, 246)
SECONDARY_GREEN = RGBColor(16, 185, 129)
GREEN_LIGHT = RGBColor(5, 150, 105)
ACCENT_ORANGE = RGBColor(245, 158, 11)
ORANGE_LIGHT = RGBColor(249, 115, 22)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(236, 72, 153)
RED = RGBColor(239, 68, 68)
DARK_TEXT = RGBColor(31, 41, 55)
LIGHT_GRAY = RGBColor(249, 250, 251)
MEDIUM_GRAY = RGBColor(229, 231, 235)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

def set_gradient_fill(shape, color1, color2, angle=45):
    """Add gradient fill to a shape"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_shadow(shape, blur=Pt(15), distance=Pt(8)):
    """Add shadow effect to shape"""
    try:
        shadow = shape.shadow
        shadow.inherit = False
        shadow.angle = 315
        shadow.blur_radius = blur
        shadow.distance = distance
        shadow.alpha = 0.5
    except:
        pass

def create_rounded_rect(slide, left, top, width, height, color, text="", font_size=20, text_color=WHITE, bold=False):
    """Create a rounded rectangle with text"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    add_shadow(shape)

    if text:
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.text = text
        p = text_frame.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER

    return shape

def add_title_with_accent(slide, title_text, y_pos=Inches(0.5)):
    """Add title with accent bar"""
    # Accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), y_pos,
        Inches(0.25), Inches(0.7)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_ORANGE
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

def slide_01_title(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Gradient background
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = PRIMARY_BLUE
    fill.gradient_stops[1].color.rgb = BLUE_LIGHT

    # Main title with glow effect
    title = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.2))
    tf = title.text_frame
    tf.text = "TradeBridge"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    tf = subtitle.text_frame
    tf.text = "Smart Supply–Demand Management System"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE

    # Institution box
    inst_box = create_rounded_rect(
        slide, Inches(2), Inches(4.8), Inches(6), Inches(1.2),
        RGBColor(100, 150, 255), font_size=16
    )
    tf = inst_box.text_frame
    tf.clear()
    tf.text = "Adama Science and Technology University\nCollege of Electrical Engineering and Computing\nDepartment of Computer Science and Engineering"
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.color.rgb = WHITE

    # Team members
    team = slide.shapes.add_textbox(Inches(1.5), Inches(6.3), Inches(7), Inches(0.5))
    tf = team.text_frame
    tf.text = "Hidaya Nurmeka • Ebisa Gutema • Hana Kebede • Hana Jote • Ilham Mohammedhassen"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE

    # Advisor
    advisor = slide.shapes.add_textbox(Inches(1.5), Inches(6.8), Inches(7), Inches(0.3))
    tf = advisor.text_frame
    tf.text = "Advisor: Dr. Ejigu Tefere"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

def slide_02_agenda(prs):
    """Slide 2: Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    add_title_with_accent(slide, "Presentation Outline")

    agenda_items = [
        "Problem Overview", "System Objectives", "Key Features",
        "System Architecture", "Technology Stack", "Implementation Highlights",
        "Machine Learning Integration", "Benefits & Impact",
        "Demo/Prototype", "Conclusion"
    ]

    # Create grid
    cols = 2
    rows = 5
    x_start = Inches(0.8)
    y_start = Inches(1.5)
    box_width = Inches(4.2)
    box_height = Inches(0.9)
    gap = Inches(0.25)

    for i, item in enumerate(agenda_items):
        row = i % rows
        col = i // rows

        x = x_start + col * (box_width + gap)
        y = y_start + row * (box_height + gap)

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x, y + Inches(0.15),
            Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_ORANGE
        circle.line.fill.background()
        add_shadow(circle, Pt(8), Pt(4))

        tf = circle.text_frame
        tf.text = str(i + 1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Text box
        text_box = slide.shapes.add_textbox(
            x + Inches(0.75), y, box_width - Inches(0.75), box_height
        )
        tf = text_box.text_frame
        tf.text = item
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT
        p.font.bold = True

def slide_03_problem(prs):
    """Slide 3: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_title_with_accent(slide, "Current Supply Chain Challenges in Ethiopia")

    problems = [
        ("⏱️", "Manual, time-consuming procurement processes"),
        ("🚫", "Lack of centralized B2B marketplace"),
        ("📉", "Poor demand planning and forecasting"),
        ("🔍", "Limited supplier visibility and comparison"),
        ("💬", "Inefficient communication between stakeholders"),
        ("📦", "Stock shortages and delivery delays")
    ]

    # Create 2x3 grid
    x_positions = [Inches(0.7), Inches(5.3)]
    y_start = Inches(1.8)
    box_width = Inches(4.2)
    box_height = Inches(1.5)
    gap = Inches(0.25)

    for i, (emoji, text) in enumerate(problems):
        row = i // 2
        col = i % 2

        x = x_positions[col]
        y = y_start + row * (box_height + gap)

        # Card
        card = create_rounded_rect(
            slide, x, y, box_width, box_height,
            WHITE, font_size=18, text_color=DARK_TEXT
        )
        card.line.color.rgb = RED
        card.line.width = Pt(4)

        # Emoji
        emoji_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(0.8), Inches(0.8))
        tf = emoji_box.text_frame
        tf.text = emoji
        p = tf.paragraphs[0]
        p.font.size = Pt(40)

        # Text
        text_box = slide.shapes.add_textbox(
            x + Inches(1.1), y + Inches(0.2), box_width - Inches(1.3), box_height - Inches(0.4)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.font.bold = True

def slide_04_intro(prs):
    """Slide 4: What is TradeBridge?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = WHITE
    fill.gradient_stops[1].color.rgb = LIGHT_GRAY

    add_title_with_accent(slide, "Introducing TradeBridge")

    # Description box
    desc_box = create_rounded_rect(
        slide, Inches(1), Inches(1.6), Inches(8), Inches(1.8),
        PRIMARY_BLUE, font_size=20, text_color=WHITE
    )
    set_gradient_fill(desc_box, PRIMARY_BLUE, BLUE_LIGHT)

    tf = desc_box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.text = "A comprehensive B2B digital platform that connects retailers, factories, distributors, and delivery personnel to streamline bulk ordering, improve supply chain transparency, and enable data-driven decision-making in the Ethiopian wholesale market."
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Stakeholder boxes
    stakeholders = [
        ("🏪", "Retailers", SECONDARY_GREEN),
        ("🏭", "Factories", PRIMARY_BLUE),
        ("🚚", "Distributors", ACCENT_ORANGE),
        ("📦", "Delivery\nPersonnel", PURPLE)
    ]

    x_start = Inches(1.5)
    y = Inches(4.2)
    box_width = Inches(1.7)
    box_height = Inches(2)
    gap = Inches(0.3)

    darker_colors = [
        RGBColor(5, 150, 105),   # Darker green
        RGBColor(29, 78, 216),   # Darker blue
        RGBColor(217, 119, 6),   # Darker orange
        RGBColor(109, 40, 217)   # Darker purple
    ]

    for i, (emoji, name, color) in enumerate(stakeholders):
        x = x_start + i * (box_width + gap)

        card = create_rounded_rect(slide, x, y, box_width, box_height, color)
        set_gradient_fill(card, color, darker_colors[i], angle=135)

        # Emoji
        emoji_box = slide.shapes.add_textbox(x, y + Inches(0.3), box_width, Inches(0.8))
        tf = emoji_box.text_frame
        tf.text = emoji
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(50)

        # Name
        name_box = slide.shapes.add_textbox(x, y + Inches(1.2), box_width, Inches(0.7))
        tf = name_box.text_frame
        tf.text = name
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE

def slide_05_scope(prs):
    """Slide 5: Project Scope"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    add_title_with_accent(slide, "Scope & Focus")

    # In Scope
    in_box = create_rounded_rect(
        slide, Inches(0.7), Inches(1.6), Inches(4.3), Inches(5),
        SECONDARY_GREEN
    )
    set_gradient_fill(in_box, SECONDARY_GREEN, GREEN_LIGHT, 135)

    # Title
    title_in = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(3.9), Inches(0.6))
    tf = title_in.text_frame
    tf.text = "✅ In Scope"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    in_items = [
        "Food and beverage products",
        "Micro to large-sized enterprises",
        "Ethiopian market (ETB currency)",
        "Mobile & Web platforms",
        "Distribution of finished goods"
    ]

    y = Inches(2.7)
    for item in in_items:
        text_box = slide.shapes.add_textbox(Inches(1), y, Inches(3.6), Inches(0.65))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = f"• {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(17)
        p.font.color.rgb = WHITE
        y += Inches(0.75)

    # Out of Scope
    out_box = create_rounded_rect(
        slide, Inches(5.3), Inches(1.6), Inches(4.3), Inches(5),
        RED
    )
    set_gradient_fill(out_box, RED, RGBColor(220, 38, 38), 135)

    # Title
    title_out = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(3.9), Inches(0.6))
    tf = title_out.text_frame
    tf.text = "❌ Out of Scope"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    out_items = [
        "Raw material procurement between factories",
        "International trade",
        "Very large national producers (e.g., Wenji Sugar)"
    ]

    y = Inches(2.7)
    for item in out_items:
        text_box = slide.shapes.add_textbox(Inches(5.6), y, Inches(3.6), Inches(0.9))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = f"• {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(17)
        p.font.color.rgb = WHITE
        y += Inches(1)

def slide_06_objectives(prs):
    """Slide 6: System Objectives"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_title_with_accent(slide, "Key Objectives")

    # General objective
    gen_box = create_rounded_rect(
        slide, Inches(1), Inches(1.5), Inches(8), Inches(1.1),
        PRIMARY_BLUE
    )
    set_gradient_fill(gen_box, PRIMARY_BLUE, BLUE_LIGHT, 90)

    tf = gen_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "General Objective: "
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Design and develop a digital platform to streamline B2B procurement and enhance supply chain visibility."
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Specific objectives
    objectives = [
        ("📱", "Develop Web & Mobile Application"),
        ("🤖", "Implement ML-based Supplier Recommendation"),
        ("📊", "Introduce Demand Forecasting Capabilities"),
        ("💬", "Enable Real-time Communication"),
        ("🔍", "Provide Centralized Supplier Directory")
    ]

    y = Inches(2.9)
    for emoji, text in objectives:
        # Card
        card = create_rounded_rect(
            slide, Inches(1.2), y, Inches(7.6), Inches(0.72),
            LIGHT_GRAY, text_color=DARK_TEXT
        )
        card.line.color.rgb = PRIMARY_BLUE
        card.line.width = Pt(2)

        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(1.5), y + Inches(0.05), Inches(0.6), Inches(0.6))
        tf = emoji_box.text_frame
        tf.text = emoji
        p = tf.paragraphs[0]
        p.font.size = Pt(32)

        # Text
        text_box = slide.shapes.add_textbox(Inches(2.3), y + Inches(0.1), Inches(6.3), Inches(0.5))
        tf = text_box.text_frame
        tf.text = text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_TEXT
        p.font.bold = True

        y += Inches(0.85)

def slide_07_features_users(prs):
    """Slide 7: Core Features - Users"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    add_title_with_accent(slide, "Platform Features - Users")

    # Retailers
    ret_box = create_rounded_rect(
        slide, Inches(0.7), Inches(1.6), Inches(4.3), Inches(5),
        PRIMARY_BLUE
    )
    set_gradient_fill(ret_box, PRIMARY_BLUE, BLUE_LIGHT, 135)

    title = slide.shapes.add_textbox(Inches(0.9), Inches(1.75), Inches(3.9), Inches(0.5))
    tf = title.text_frame
    tf.text = "🏪 For Retailers"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    items = [
        "Browse and compare products from multiple suppliers",
        "Place bulk orders with cart management",
        "Track order status in real-time",
        "Rate and review suppliers",
        "Receive personalized supplier recommendations"
    ]

    y = Inches(2.5)
    for item in items:
        text_box = slide.shapes.add_textbox(Inches(1), y, Inches(3.6), Inches(0.75))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = f"• {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        y += Inches(0.85)

    # Distributors/Factories
    dist_box = create_rounded_rect(
        slide, Inches(5.3), Inches(1.6), Inches(4.3), Inches(5),
        ACCENT_ORANGE
    )
    set_gradient_fill(dist_box, ACCENT_ORANGE, ORANGE_LIGHT, 135)

    title = slide.shapes.add_textbox(Inches(5.5), Inches(1.75), Inches(3.9), Inches(0.5))
    tf = title.text_frame
    tf.text = "🏭 For Distributors/Factories"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    items2 = [
        "Manage product listings and inventory",
        "Approve/reject incoming orders",
        "Broadcast promotional announcements",
        "View demand analytics and sales reports"
    ]

    y = Inches(2.5)
    for item in items2:
        text_box = slide.shapes.add_textbox(Inches(5.6), y, Inches(3.6), Inches(0.9))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = f"• {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        y += Inches(1)

def slide_08_features_smart(prs):
    """Slide 8: Core Features - Smart Capabilities"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_title_with_accent(slide, "Platform Features - Smart Capabilities")

    # ML Features
    ml_box1 = create_rounded_rect(
        slide, Inches(1), Inches(1.6), Inches(8), Inches(1.3),
        PURPLE
    )
    set_gradient_fill(ml_box1, PURPLE, RGBColor(124, 58, 237), 90)

    tf = ml_box1.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "🤖 Supplier Recommendation System"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Ranks suppliers based on price, distance, performance, and user preferences.\nPersonalized recommendations for each retailer."
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)

    ml_box2 = create_rounded_rect(
        slide, Inches(1), Inches(3.1), Inches(8), Inches(1.3),
        SECONDARY_GREEN
    )
    set_gradient_fill(ml_box2, SECONDARY_GREEN, GREEN_LIGHT, 90)

    tf = ml_box2.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "📈 Demand Forecasting"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Predicts future product demand for manufacturers.\nReduces stockouts and overstock situations."
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)

    # Other features
    features = [
        ("📍", "Real-time GPS Tracking", PRIMARY_BLUE),
        ("💳", "Chapa Payment", PINK),
        ("📧", "Auto Notifications", ACCENT_ORANGE),
        ("💬", "In-App Chat", SECONDARY_GREEN)
    ]

    x_start = Inches(1.2)
    y = Inches(4.8)
    box_width = Inches(1.9)
    box_height = Inches(1.6)
    gap = Inches(0.15)

    feature_darker_colors = [
        RGBColor(29, 78, 216),   # Darker blue
        RGBColor(219, 39, 119),  # Darker pink
        RGBColor(217, 119, 6),   # Darker orange
        RGBColor(5, 150, 105)    # Darker green
    ]

    for i, (emoji, name, color) in enumerate(features):
        x = x_start + i * (box_width + gap)

        card = create_rounded_rect(slide, x, y, box_width, box_height, color)
        set_gradient_fill(card, color, feature_darker_colors[i], 135)

        # Emoji
        emoji_box = slide.shapes.add_textbox(x, y + Inches(0.2), box_width, Inches(0.7))
        tf = emoji_box.text_frame
        tf.text = emoji
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(40)

        # Name
        name_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.95), box_width - Inches(0.2), Inches(0.6))
        tf = name_box.text_frame
        tf.word_wrap = True
        tf.text = name
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

def slide_09_architecture(prs):
    """Slide 9: System Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = WHITE
    fill.gradient_stops[1].color.rgb = LIGHT_GRAY

    add_title_with_accent(slide, "Three-Tier Architecture")

    layers = [
        ("Presentation Layer",
         ["Mobile App (Android)", "Web Application", "User Interface Components"],
         PRIMARY_BLUE, BLUE_LIGHT, Inches(1.7)),
        ("Application Layer (Node.js + Express)",
         ["Authentication & Authorization", "Business Logic", "API Services", "ML Model Integration"],
         SECONDARY_GREEN, GREEN_LIGHT, Inches(3.4)),
        ("Data Layer (MySQL)",
         ["User Data", "Products & Orders", "Analytics & Logs"],
         ACCENT_ORANGE, ORANGE_LIGHT, Inches(5.5))
    ]

    for title, items, color1, color2, y_pos in layers:
        # Layer box
        layer_box = create_rounded_rect(
            slide, Inches(1.5), y_pos, Inches(7), Inches(1.5),
            color1
        )
        set_gradient_fill(layer_box, color1, color2, 90)

        # Title
        title_box = slide.shapes.add_textbox(Inches(1.7), y_pos + Inches(0.15), Inches(6.6), Inches(0.4))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Items
        items_text = " • ".join(items)
        items_box = slide.shapes.add_textbox(Inches(1.7), y_pos + Inches(0.65), Inches(6.6), Inches(0.8))
        tf = items_box.text_frame
        tf.word_wrap = True
        tf.text = items_text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE

    # Arrows
    for y in [Inches(3.3), Inches(5.1)]:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            Inches(4.5), y,
            Inches(1), Inches(0.5)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_TEXT
        arrow.line.fill.background()

def slide_10_tech_stack(prs):
    """Slide 10: Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    add_title_with_accent(slide, "Technologies Used")

    tech_categories = [
        ("Frontend", ["⚛️ React.js with TypeScript", "🎨 Tailwind CSS", "📊 Zustand"],
         PRIMARY_BLUE, BLUE_LIGHT, Inches(1.7)),
        ("Backend", ["🟢 Node.js with Express", "🔐 JWT Authentication", "📡 RESTful APIs"],
         SECONDARY_GREEN, GREEN_LIGHT, Inches(2.7)),
        ("Database", ["🗄️ MySQL with Sequelize ORM"],
         ACCENT_ORANGE, ORANGE_LIGHT, Inches(3.7)),
        ("Machine Learning", ["🐍 Python", "📚 Scikit-learn, Pandas, NumPy"],
         PURPLE, RGBColor(124, 58, 237), Inches(4.7)),
        ("Payment", ["💰 Chapa Payment Gateway"],
         PINK, RGBColor(219, 39, 119), Inches(5.7))
    ]

    for category, items, color1, color2, y_pos in tech_categories:
        tech_box = create_rounded_rect(
            slide, Inches(1), y_pos, Inches(8), Inches(0.8),
            color1
        )
        set_gradient_fill(tech_box, color1, color2, 90)

        # Category name
        cat_box = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.1), Inches(2), Inches(0.6))
        tf = cat_box.text_frame
        tf.text = category
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Items
        items_text = " • ".join(items)
        items_box = slide.shapes.add_textbox(Inches(3.3), y_pos + Inches(0.1), Inches(5.5), Inches(0.6))
        tf = items_box.text_frame
        tf.word_wrap = True
        tf.text = items_text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(17)
        p.font.color.rgb = WHITE

def slide_11_user_roles(prs):
    """Slide 11: User Roles & Access"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    add_title_with_accent(slide, "Multi-Role Access Control")

    roles = [
        ("👤", "Retailer", "Browse products, Place orders, Track deliveries, Rate suppliers", PRIMARY_BLUE),
        ("🏭", "Factory", "Manage products, Approve orders, View demand forecasts", SECONDARY_GREEN),
        ("🚚", "Distributor", "Buy & sell, Manage inventory, Fulfill orders", ACCENT_ORANGE),
        ("🚗", "Driver", "View assignments, Update delivery status, Track routes", PURPLE),
        ("👑", "Admin", "Manage users, Approve suppliers, Monitor platform", PINK)
    ]

    y = Inches(1.7)
    for emoji, role, permissions, color in roles:
        # Role card
        card = create_rounded_rect(
            slide, Inches(1), y, Inches(8), Inches(0.85),
            WHITE
        )
        card.line.color.rgb = color
        card.line.width = Pt(3)

        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(1.3), y + Inches(0.1), Inches(0.6), Inches(0.6))
        tf = emoji_box.text_frame
        tf.text = emoji
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(36)

        # Role name
        role_box = slide.shapes.add_textbox(Inches(2.1), y + Inches(0.15), Inches(1.5), Inches(0.5))
        tf = role_box.text_frame
        tf.text = role
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = color

        # Permissions
        perm_box = slide.shapes.add_textbox(Inches(3.8), y + Inches(0.15), Inches(4.9), Inches(0.55))
        tf = perm_box.text_frame
        tf.word_wrap = True
        tf.text = permissions
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT

        y += Inches(1)

def slide_12_database(prs):
    """Slide 12: Database Design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_title_with_accent(slide, "Core Data Entities")

    entities = [
        ("👥", "Users", "retailers, suppliers,\ndrivers, admins"),
        ("📦", "Products", "name, price, stock,\nMOQ, images"),
        ("🛒", "Orders", "status, items,\ntracking, payments"),
        ("💬", "Messages", "in-app chat\nhistory"),
        ("⭐", "Ratings &\nReviews", "supplier performance\nmetrics"),
        ("💳", "Payments", "transactions, methods,\nstatus")
    ]

    # Create 3x2 grid
    x_positions = [Inches(0.8), Inches(3.8), Inches(6.8)]
    y_positions = [Inches(1.8), Inches(4.5)]

    for i, (emoji, name, desc) in enumerate(entities):
        row = i // 3
        col = i % 3

        x = x_positions[col]
        y = y_positions[row]

        # Entity card
        card = create_rounded_rect(
            slide, x, y, Inches(2.7), Inches(2.3),
            PRIMARY_BLUE
        )
        set_gradient_fill(card, PRIMARY_BLUE, BLUE_LIGHT, 135)

        # Emoji
        emoji_box = slide.shapes.add_textbox(x, y + Inches(0.2), Inches(2.7), Inches(0.6))
        tf = emoji_box.text_frame
        tf.text = emoji
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(50)

        # Name
        name_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.9), Inches(2.5), Inches(0.5))
        tf = name_box.text_frame
        tf.word_wrap = True
        tf.text = name
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Description
        desc_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.5), Inches(2.5), Inches(0.7))
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = desc
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE

def slide_13_ml_recommendation(prs):
    """Slide 13: ML Supplier Recommendation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    add_title_with_accent(slide, "Intelligent Supplier Ranking")

    # Process steps
    steps = [
        ("1", "Collects Data", "Price, delivery time,\nratings, location", PRIMARY_BLUE),
        ("2", "Trains Model", "Random Forest\nClassifier", SECONDARY_GREEN),
        ("3", "Generates Score", "Ranks suppliers\nfor each retailer", ACCENT_ORANGE),
        ("4", "Personalizes", "Based on past\npreferences", PURPLE)
    ]

    x_start = Inches(0.9)
    y = Inches(1.7)
    box_width = Inches(2.1)
    box_height = Inches(1.8)
    gap = Inches(0.2)

    for i, (num, title, desc, color) in enumerate(steps):
        x = x_start + i * (box_width + gap)

        # Step card
        darker_step_colors = [
            RGBColor(29, 78, 216),   # Darker blue
            RGBColor(5, 150, 105),   # Darker green
            RGBColor(217, 119, 6),   # Darker orange
            RGBColor(109, 40, 217)   # Darker purple
        ]
        card = create_rounded_rect(slide, x, y, box_width, box_height, color)
        set_gradient_fill(card, color, darker_step_colors[i], 135)

        # Number
        num_box = slide.shapes.add_textbox(x, y + Inches(0.15), box_width, Inches(0.5))
        tf = num_box.text_frame
        tf.text = num
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Title
        title_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.75), box_width - Inches(0.2), Inches(0.4))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Description
        desc_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.2), box_width - Inches(0.2), Inches(0.55))
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = desc
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE

        # Arrow
        if i < 3:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + box_width + Inches(0.02), y + Inches(0.75),
                Inches(0.16), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_TEXT
            arrow.line.fill.background()

    # Features box
    features_box = create_rounded_rect(
        slide, Inches(1), Inches(3.8), Inches(8), Inches(2.6),
        WHITE
    )
    features_box.line.color.rgb = PRIMARY_BLUE
    features_box.line.width = Pt(3)

    title = slide.shapes.add_textbox(Inches(1.3), Inches(4), Inches(7.4), Inches(0.5))
    tf = title.text_frame
    tf.text = "Features Used for Ranking:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    features = [
        "✓ Price competitiveness",
        "✓ On-time delivery rate",
        "✓ Quality ratings",
        "✓ Fulfillment time",
        "✓ Communication responsiveness"
    ]

    y_feat = Inches(4.7)
    for feature in features:
        feat_box = slide.shapes.add_textbox(Inches(1.5), y_feat, Inches(7), Inches(0.35))
        tf = feat_box.text_frame
        tf.text = feature
        p = tf.paragraphs[0]
        p.font.size = Pt(19)
        p.font.color.rgb = DARK_TEXT
        y_feat += Inches(0.38)

def slide_14_ml_forecasting(prs):
    """Slide 14: ML Demand Forecasting"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_title_with_accent(slide, "Predictive Demand Planning")

    # Purpose box
    purpose_box = create_rounded_rect(
        slide, Inches(1), Inches(1.6), Inches(8), Inches(1.1),
        PURPLE
    )
    set_gradient_fill(purpose_box, PURPLE, RGBColor(124, 58, 237), 90)

    tf = purpose_box.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "Purpose: "
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Help manufacturers plan production based on predicted future demand"
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    # Approach box
    approach_box = create_rounded_rect(
        slide, Inches(0.7), Inches(3), Inches(4.3), Inches(3.4),
        PRIMARY_BLUE
    )
    set_gradient_fill(approach_box, PRIMARY_BLUE, BLUE_LIGHT, 135)

    title = slide.shapes.add_textbox(Inches(0.9), Inches(3.15), Inches(3.9), Inches(0.5))
    tf = title.text_frame
    tf.text = "Approach"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    items_approach = [
        "Algorithm: Linear Regression,\nRandom Forest Regressor",
        "Features: Historical demand,\nmoving averages, seasonality",
        "Output: Predicted demand\nquantities",
        "Evaluation: MAE, RMSE"
    ]

    y = Inches(3.8)
    for item in items_approach:
        text_box = slide.shapes.add_textbox(Inches(1), y, Inches(3.6), Inches(0.6))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = f"• {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        y += Inches(0.68)

    # Benefits box
    benefits_box = create_rounded_rect(
        slide, Inches(5.3), Inches(3), Inches(4.3), Inches(3.4),
        SECONDARY_GREEN
    )
    set_gradient_fill(benefits_box, SECONDARY_GREEN, GREEN_LIGHT, 135)

    title = slide.shapes.add_textbox(Inches(5.5), Inches(3.15), Inches(3.9), Inches(0.5))
    tf = title.text_frame
    tf.text = "Benefits"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    items_benefits = [
        "✅ Reduce stockouts",
        "✅ Minimize overproduction",
        "✅ Optimize inventory\nmanagement",
        "✅ Data-driven production\nplanning"
    ]

    y = Inches(3.8)
    for item in items_benefits:
        text_box = slide.shapes.add_textbox(Inches(5.6), y, Inches(3.6), Inches(0.6))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(17)
        p.font.color.rgb = WHITE
        y += Inches(0.7)

def slide_15_payment(prs):
    """Slide 15: Payment Integration"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    add_title_with_accent(slide, "Secure Payment Processing")

    # Gateway banner
    gateway_box = create_rounded_rect(
        slide, Inches(2.5), Inches(1.6), Inches(5), Inches(0.9),
        PINK
    )
    set_gradient_fill(gateway_box, PINK, RGBColor(219, 39, 119), 90)

    tf = gateway_box.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "💰 Chapa Payment Gateway"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Compliance
    comp_box = create_rounded_rect(
        slide, Inches(0.7), Inches(2.8), Inches(4.3), Inches(3.6),
        WHITE
    )
    comp_box.line.color.rgb = PRIMARY_BLUE
    comp_box.line.width = Pt(3)

    title = slide.shapes.add_textbox(Inches(0.9), Inches(2.95), Inches(3.9), Inches(0.5))
    tf = title.text_frame
    tf.text = "✅ Compliance"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    comp_items = [
        "NBE Directive ONPS/10/2025\n(Transaction limits)",
        "KYC/AML verification",
        "Split-payment mechanism",
        "Secure record keeping"
    ]

    y = Inches(3.6)
    for item in comp_items:
        text_box = slide.shapes.add_textbox(Inches(1), y, Inches(3.6), Inches(0.65))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = f"• {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        y += Inches(0.75)

    # Methods & Security
    methods_box = create_rounded_rect(
        slide, Inches(5.3), Inches(2.8), Inches(4.3), Inches(3.6),
        WHITE
    )
    methods_box.line.color.rgb = SECONDARY_GREEN
    methods_box.line.width = Pt(3)

    title = slide.shapes.add_textbox(Inches(5.5), Inches(2.95), Inches(3.9), Inches(0.5))
    tf = title.text_frame
    tf.text = "💳 Supported Methods"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_GREEN

    methods = ["Mobile wallets", "E-money", "Bank transfers"]
    y = Inches(3.6)
    for item in methods:
        text_box = slide.shapes.add_textbox(Inches(5.6), y, Inches(3.6), Inches(0.4))
        tf = text_box.text_frame
        tf.text = f"• {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(17)
        p.font.color.rgb = DARK_TEXT
        y += Inches(0.45)

    # Security
    y = Inches(4.9)
    title = slide.shapes.add_textbox(Inches(5.5), y, Inches(3.9), Inches(0.4))
    tf = title.text_frame
    tf.text = "🔐 Security"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE

    security = ["Encrypted storage", "JWT authentication", "Transaction logging"]
    y = Inches(5.4)
    for item in security:
        text_box = slide.shapes.add_textbox(Inches(5.6), y, Inches(3.6), Inches(0.4))
        tf = text_box.text_frame
        tf.text = f"• {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        y += Inches(0.43)

def slide_16_benefits(prs):
    """Slide 16: Benefits & Impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_title_with_accent(slide, "Expected Impact")

    # Three columns
    columns = [
        ("For Retailers",
         ["⚡ Faster procurement process", "💰 Better price comparison",
          "📊 Improved supplier visibility", "🎯 Personalized recommendations"],
         PRIMARY_BLUE, BLUE_LIGHT, Inches(0.7)),
        ("For Suppliers",
         ["📈 Expanded market reach", "🤝 Direct buyer connections",
          "📉 Reduced manual operations", "📊 Access to demand insights"],
         SECONDARY_GREEN, GREEN_LIGHT, Inches(3.7)),
        ("For the Industry",
         ["🌐 Digital transformation", "📉 Reduced inefficiencies",
          "🔍 Increased transparency", "📈 Data-driven decisions"],
         ACCENT_ORANGE, ORANGE_LIGHT, Inches(6.7))
    ]

    for title, items, color1, color2, x_pos in columns:
        # Column box
        col_box = create_rounded_rect(
            slide, x_pos, Inches(1.6), Inches(2.9), Inches(4.8),
            WHITE
        )
        col_box.line.color.rgb = color1
        col_box.line.width = Pt(4)

        # Header
        header = create_rounded_rect(
            slide, x_pos, Inches(1.6), Inches(2.9), Inches(0.7),
            color1
        )
        set_gradient_fill(header, color1, color2, 90)

        tf = header.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Items
        y = Inches(2.5)
        for item in items:
            item_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y, Inches(2.5), Inches(0.7))
            tf = item_box.text_frame
            tf.word_wrap = True
            tf.text = item
            p = tf.paragraphs[0]
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_TEXT
            y += Inches(0.88)

def slide_17_feasibility(prs):
    """Slide 17: Feasibility Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    add_title_with_accent(slide, "Project Feasibility")

    feasibility_data = [
        ("Technical Feasibility",
         ["Proven technologies\n(React, Node.js, MySQL)", "Team has required skills",
          "No specialized hardware"],
         PRIMARY_BLUE, Inches(0.8)),
        ("Operational Feasibility",
         ["User-friendly interface", "Supports existing\nbusiness practices",
          "Mobile-first approach"],
         SECONDARY_GREEN, Inches(3.7)),
        ("Economic Feasibility",
         ["Low development cost\n(~9,500 ETB)", "Open-source technologies",
          "Revenue through fees & ads"],
         ACCENT_ORANGE, Inches(6.6))
    ]

    for title, items, color, x_pos in feasibility_data:
        # Card
        card = create_rounded_rect(
            slide, x_pos, Inches(1.8), Inches(2.9), Inches(4.6),
            WHITE
        )
        card.line.color.rgb = color
        card.line.width = Pt(4)

        # Checkmark
        check_box = slide.shapes.add_textbox(x_pos, Inches(2), Inches(2.9), Inches(0.8))
        tf = check_box.text_frame
        tf.text = "✅"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(60)

        # Title
        title_box = slide.shapes.add_textbox(x_pos + Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = title
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(19)
        p.font.bold = True
        p.font.color.rgb = color

        # Items
        y = Inches(3.7)
        for item in items:
            item_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y, Inches(2.5), Inches(0.7))
            tf = item_box.text_frame
            tf.word_wrap = True
            tf.text = f"• {item}"
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_TEXT
            y += Inches(0.85)

def slide_18_methodology(prs):
    """Slide 18: Development Methodology"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_title_with_accent(slide, "Agile & Incremental Approach")

    phases = [
        ("📋", "Planning", "Requirements gathering\n(Week 1-2)", PRIMARY_BLUE),
        ("🔍", "Analysis", "UML design, data modeling\n(Week 3-8)", SECONDARY_GREEN),
        ("🎨", "Design", "Architecture, UI/UX\n(Week 8-12)", ACCENT_ORANGE),
        ("💻", "Implementation", "Frontend, backend, ML\n(Week 13-15)", PURPLE),
        ("🧪", "Testing", "System validation\n(Week 15-16)", PINK),
        ("📄", "Documentation", "Final reports\n(Week 17-18)", DARK_TEXT)
    ]

    # Create 3x2 grid
    x_positions = [Inches(0.8), Inches(3.8), Inches(6.8)]
    y_positions = [Inches(1.8), Inches(4.2)]

    for i, (emoji, phase, desc, color) in enumerate(phases):
        row = i // 3
        col = i % 3

        x = x_positions[col]
        y = y_positions[row]

        # Phase card
        card = create_rounded_rect(slide, x, y, Inches(2.7), Inches(2), color)

        # Emoji
        emoji_box = slide.shapes.add_textbox(x, y + Inches(0.15), Inches(2.7), Inches(0.5))
        tf = emoji_box.text_frame
        tf.text = emoji
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(42)

        # Phase name
        phase_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.75), Inches(2.5), Inches(0.4))
        tf = phase_box.text_frame
        tf.word_wrap = True
        tf.text = phase
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Description
        desc_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.25), Inches(2.5), Inches(0.7))
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = desc
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE

    # Timeline banner
    timeline_box = create_rounded_rect(
        slide, Inches(2.5), Inches(6.5), Inches(5), Inches(0.7),
        SECONDARY_GREEN
    )
    set_gradient_fill(timeline_box, SECONDARY_GREEN, GREEN_LIGHT, 90)

    tf = timeline_box.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "Timeline: ~18 weeks (2 semesters)"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def slide_19_challenges(prs):
    """Slide 19: Challenges & Limitations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    add_title_with_accent(slide, "Acknowledged Limitations")

    challenges = [
        ("📊", "Limited initial ML training data"),
        ("🌐", "Dependency on stable internet connectivity"),
        ("🤔", "Resistance to digital adoption from traditional businesses"),
        ("💳", "Reliance on third-party payment providers"),
        ("📚", "User training needs")
    ]

    y = Inches(1.7)
    for emoji, text in challenges:
        # Challenge card
        card = create_rounded_rect(
            slide, Inches(1), y, Inches(8), Inches(0.7),
            WHITE
        )
        card.line.color.rgb = RED
        card.line.width = Pt(3)

        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(1.3), y + Inches(0.05), Inches(0.6), Inches(0.6))
        tf = emoji_box.text_frame
        tf.text = emoji
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(32)

        # Text
        text_box = slide.shapes.add_textbox(Inches(2.1), y + Inches(0.1), Inches(6.6), Inches(0.5))
        tf = text_box.text_frame
        tf.text = text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT

        y += Inches(0.8)

    # Mitigation
    mitigation_box = create_rounded_rect(
        slide, Inches(1), Inches(5.4), Inches(8), Inches(1.9),
        SECONDARY_GREEN
    )
    set_gradient_fill(mitigation_box, SECONDARY_GREEN, GREEN_LIGHT, 90)

    title = slide.shapes.add_textbox(Inches(1.2), Inches(5.55), Inches(7.6), Inches(0.5))
    tf = title.text_frame
    tf.text = "🛡️ Mitigation Strategies"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    mitigations = [
        "Hybrid ML approach (rule-based + data-driven)",
        "Offline capabilities for critical functions",
        "User onboarding and support system"
    ]

    y = Inches(6.2)
    for item in mitigations:
        item_box = slide.shapes.add_textbox(Inches(1.5), y, Inches(7), Inches(0.35))
        tf = item_box.text_frame
        tf.text = f"✓ {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        y += Inches(0.4)

def slide_20_future(prs):
    """Slide 20: Future Enhancements"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_title_with_accent(slide, "Roadmap & Future Work")

    enhancements = [
        ("🌍", "Support for additional product categories"),
        ("🤖", "Advanced AI-powered chatbot support"),
        ("📱", "iOS mobile application"),
        ("🌐", "Multi-language support (Amharic, Oromifa, etc.)"),
        ("📊", "Advanced analytics dashboards"),
        ("🔗", "Integration with ERP systems"),
        ("🚛", "Advanced logistics optimization"),
        ("🔐", "Blockchain for supply chain transparency")
    ]

    # Create 2x4 grid
    x_positions = [Inches(0.8), Inches(5.4)]
    y_start = Inches(1.7)
    box_width = Inches(4.2)
    box_height = Inches(0.95)
    gap = Inches(0.2)

    for i, (emoji, text) in enumerate(enhancements):
        row = i // 2
        col = i % 2

        x = x_positions[col]
        y = y_start + row * (box_height + gap)

        # Enhancement card
        card = create_rounded_rect(
            slide, x, y, box_width, box_height,
            LIGHT_GRAY
        )
        card.line.color.rgb = PRIMARY_BLUE
        card.line.width = Pt(2)

        # Emoji
        emoji_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.6))
        tf = emoji_box.text_frame
        tf.text = emoji
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(32)

        # Text
        text_box = slide.shapes.add_textbox(x + Inches(0.9), y + Inches(0.15), box_width - Inches(1.1), Inches(0.65))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.font.bold = True

def slide_21_conclusion(prs):
    """Slide 21: Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = PRIMARY_BLUE
    fill.gradient_stops[1].color.rgb = BLUE_LIGHT

    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(8), Inches(0.7))
    tf = title.text_frame
    tf.text = "Summary"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Deliverables box
    deliv_box = create_rounded_rect(
        slide, Inches(1.2), Inches(1.7), Inches(7.6), Inches(3.2),
        WHITE
    )

    title_deliv = slide.shapes.add_textbox(Inches(1.4), Inches(1.85), Inches(7.2), Inches(0.5))
    tf = title_deliv.text_frame
    tf.text = "TradeBridge delivers:"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    deliverables = [
        "✅ Centralized B2B marketplace for Ethiopian wholesale",
        "✅ Smart supplier recommendations & demand forecasting",
        "✅ Real-time tracking & secure payments",
        "✅ Improved efficiency across the supply chain"
    ]

    y = Inches(2.5)
    for item in deliverables:
        item_box = slide.shapes.add_textbox(Inches(1.7), y, Inches(6.8), Inches(0.35))
        tf = item_box.text_frame
        tf.word_wrap = True
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        y += Inches(0.5)

    # Outcomes
    y = Inches(4.1)
    title_out = slide.shapes.add_textbox(Inches(1.4), y, Inches(7.2), Inches(0.4))
    tf = title_out.text_frame
    tf.text = "Expected Outcomes:"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_GREEN

    outcomes = [
        "Digital transformation of wholesale procurement",
        "Reduced operational costs",
        "Enhanced transparency and trust"
    ]

    y = Inches(4.6)
    for item in outcomes:
        item_box = slide.shapes.add_textbox(Inches(1.7), y, Inches(6.8), Inches(0.25))
        tf = item_box.text_frame
        tf.text = f"• {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        y += Inches(0.3)

    # Status banner
    status_box = create_rounded_rect(
        slide, Inches(1.8), Inches(5.6), Inches(6.4), Inches(0.8),
        SECONDARY_GREEN
    )
    set_gradient_fill(status_box, SECONDARY_GREEN, GREEN_LIGHT, 90)

    tf = status_box.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "🚀 Status: Prototype completed, ready for deployment"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def slide_22_thankyou(prs):
    """Slide 22: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = SECONDARY_GREEN
    fill.gradient_stops[1].color.rgb = GREEN_LIGHT

    # Thank you
    thanks = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
    tf = thanks.text_frame
    tf.text = "Thank You!"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Questions
    questions = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    tf = questions.text_frame
    tf.text = "Questions?"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE

    # Contact info
    contact = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(6), Inches(0.9))
    tf = contact.text_frame
    tf.text = "Advisor: Dr. Ejigu Tefere\nAdama Science and Technology University"
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(20)
        para.font.color.rgb = WHITE

    # Team
    team = slide.shapes.add_textbox(Inches(1.5), Inches(6), Inches(7), Inches(0.6))
    tf = team.text_frame
    tf.word_wrap = True
    tf.text = "Hidaya Nurmeka • Ebisa Gutema • Hana Kebede • Hana Jote • Ilham Mohammedhassen"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(15)
    p.font.color.rgb = WHITE

def create_all_slides(prs):
    """Create all 22 slides"""
    print("Creating slide 1: Title")
    slide_01_title(prs)

    print("Creating slide 2: Agenda")
    slide_02_agenda(prs)

    print("Creating slide 3: Problem")
    slide_03_problem(prs)

    print("Creating slide 4: Introduction")
    slide_04_intro(prs)

    print("Creating slide 5: Scope")
    slide_05_scope(prs)

    print("Creating slide 6: Objectives")
    slide_06_objectives(prs)

    print("Creating slide 7: Features - Users")
    slide_07_features_users(prs)

    print("Creating slide 8: Features - Smart")
    slide_08_features_smart(prs)

    print("Creating slide 9: Architecture")
    slide_09_architecture(prs)

    print("Creating slide 10: Technology Stack")
    slide_10_tech_stack(prs)

    print("Creating slide 11: User Roles")
    slide_11_user_roles(prs)

    print("Creating slide 12: Database Design")
    slide_12_database(prs)

    print("Creating slide 13: ML Recommendation")
    slide_13_ml_recommendation(prs)

    print("Creating slide 14: ML Forecasting")
    slide_14_ml_forecasting(prs)

    print("Creating slide 15: Payment Integration")
    slide_15_payment(prs)

    print("Creating slide 16: Benefits & Impact")
    slide_16_benefits(prs)

    print("Creating slide 17: Feasibility")
    slide_17_feasibility(prs)

    print("Creating slide 18: Methodology")
    slide_18_methodology(prs)

    print("Creating slide 19: Challenges")
    slide_19_challenges(prs)

    print("Creating slide 20: Future Enhancements")
    slide_20_future(prs)

    print("Creating slide 21: Conclusion")
    slide_21_conclusion(prs)

    print("Creating slide 22: Thank You")
    slide_22_thankyou(prs)

def main():
    """Main function"""
    import os

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("\n🎨 Creating TradeBridge Premium Presentation...\n")

    create_all_slides(prs)

    output_path = 'TradeBridge_Premium_Presentation.pptx'
    prs.save(output_path)

    print(f"\n✅ Premium presentation created successfully!")
    print(f"📁 Location: {os.path.abspath(output_path)}")
    print(f"📊 Total slides: {len(prs.slides)}")

    file_size = os.path.getsize(output_path) / (1024 * 1024)
    print(f"💾 File size: {file_size:.2f} MB\n")

if __name__ == '__main__':
    main()
